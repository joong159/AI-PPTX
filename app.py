import os
import re
import json
import io
import uuid
import streamlit as st
import streamlit.components.v1 as components
from dotenv import load_dotenv
import google.generativeai as genai
from groq import Groq
from pypdf import PdfReader
from ppt_compiler import compile_presentation, THEMES, TEMPLATE_STYLES
from pptx.dml.color import RGBColor
try:
    import requests
    from bs4 import BeautifulSoup
    _WEB_FETCH_OK = True
except ImportError:
    _WEB_FETCH_OK = False

load_dotenv()


def extract_pptx_text(file_obj) -> str:
    """Upload된 .pptx 파일에서 슬라이드 텍스트를 추출합니다."""
    from pptx import Presentation
    prs = Presentation(file_obj)
    parts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            parts.append(f"[슬라이드 {i+1}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def fetch_url_text(url: str) -> tuple:
    """URL에서 텍스트 콘텐츠를 가져옵니다. (text, error) 반환."""
    if not _WEB_FETCH_OK:
        return None, "requests / beautifulsoup4 패키지가 필요합니다."
    try:
        resp = requests.get(url.strip(), timeout=12,
                            headers={"User-Agent": "Mozilla/5.0 (compatible; SnapdeckBot/1.0)"})
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, 'html.parser')
        for tag in soup(['script', 'style', 'nav', 'footer', 'header', 'aside', 'noscript']):
            tag.decompose()
        text = soup.get_text(separator='\n', strip=True)
        # 빈 줄 정리
        lines = [l for l in text.splitlines() if l.strip()]
        return '\n'.join(lines[:400]), None  # 약 5000자 상한
    except Exception as e:
        return None, str(e)


import os as _os
_COMP_DIR = _os.path.join(_os.path.dirname(__file__), "slide_editor")
_slide_editor_component = components.declare_component("slide_editor", path=_COMP_DIR)

st.set_page_config(
    page_title="Snapdeck — Build a winning deck in a snap",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ── CSS ──────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@400;500;700;800;900&display=swap');

/* ── 기본: 라이트 모드 ── */
:root {
    --bg: #F8F9FF;
    --surface: #FFFFFF;
    --elevated: #EEF0FF;
    --border: rgba(99,102,241,0.1);
    --border-b: rgba(99,102,241,0.22);
    --text: #0F0F2A;
    --text2: #4A4A6A;
    --text3: #9090B0;
    --accent: #7C3AED;
    --accent2: #4F46E5;
    --accent-l: #6D28D9;
    --accent-bg: rgba(124,58,237,0.07);
    --accent-glow: rgba(124,58,237,0.18);
    --green: #16A34A;
    --green-bg: rgba(22,163,74,0.08);
    --nav-bg: rgba(248,249,255,0.92);
    --feat-bg: #FFFFFF;
    --r-sm: 8px;
    --r-md: 12px;
    --r-lg: 20px;
    --t: 0.2s ease;
}

/* ── Base ── */
html, body,
[data-testid="stAppViewContainer"],
[data-testid="stMain"], section.main {
    background: var(--bg) !important;
    color: var(--text) !important;
}
.main .block-container {
    padding-top: 0 !important;
    padding-bottom: 4rem !important;
    max-width: 1200px !important;
    font-family: 'Noto Sans KR', 'Inter', -apple-system, BlinkMacSystemFont, sans-serif !important;
}
#MainMenu, footer, [data-testid="stHeader"] { display: none !important; }

/* ── Sidebar ── */
[data-testid="stSidebar"] {
    background: var(--surface) !important;
    border-right: 1px solid var(--border) !important;
}
[data-testid="stSidebar"] .stButton > button {
    background: var(--elevated) !important;
    color: var(--text2) !important;
    border: 1px solid var(--border) !important;
    box-shadow: none !important;
    transform: none !important;
}
[data-testid="stSidebar"] .stButton > button:hover {
    background: var(--accent-bg) !important;
    color: var(--accent-l) !important;
    border-color: var(--accent) !important;
    box-shadow: none !important;
    transform: none !important;
}

/* ── Inputs ── */
.stTextInput > div > div > input,
.stTextArea > div > textarea {
    background: var(--elevated) !important;
    color: var(--text) !important;
    border: 1px solid var(--border) !important;
    border-radius: var(--r-sm) !important;
    font-family: inherit !important;
    font-size: 0.93rem !important;
    line-height: 1.55 !important;
    transition: border-color var(--t), box-shadow var(--t) !important;
}
.stTextInput > div > div > input:focus,
.stTextArea > div > textarea:focus {
    border-color: var(--accent) !important;
    box-shadow: 0 0 0 3px var(--accent-glow) !important;
    outline: none !important;
}
.stTextInput > div > div > input::placeholder,
.stTextArea > div > textarea::placeholder { color: var(--text3) !important; }

.stSelectbox > div > div,
.stMultiSelect > div > div {
    background: var(--elevated) !important;
    color: var(--text) !important;
    border: 1px solid var(--border) !important;
    border-radius: var(--r-sm) !important;
}

/* ── Labels ── */
.stTextInput > label, .stTextArea > label,
.stSelectbox > label, .stRadio > label,
.stFileUploader > label, .stSlider > label,
.stNumberInput > label {
    color: var(--text2) !important;
    font-weight: 500 !important;
    font-size: 0.85rem !important;
    font-family: inherit !important;
    margin-bottom: 4px !important;
}
.stRadio > div > label, .stCheckbox > label {
    color: var(--text2) !important;
    font-size: 0.9rem !important;
}

/* ── Buttons ── */
.stButton > button {
    background: var(--accent) !important;
    color: #fff !important;
    border: none !important;
    border-radius: var(--r-sm) !important;
    font-weight: 600 !important;
    font-family: inherit !important;
    font-size: 0.88rem !important;
    padding: 0.55rem 1.1rem !important;
    transition: all var(--t) !important;
    letter-spacing: 0.01em !important;
}
.stButton > button:hover {
    background: #6D28D9 !important;
    transform: translateY(-1px) !important;
    box-shadow: 0 4px 14px var(--accent-glow) !important;
}
.stButton > button:active { transform: translateY(0) !important; }

.stDownloadButton > button {
    background: linear-gradient(135deg, var(--accent), #4F46E5) !important;
    color: #fff !important;
    border: none !important;
    border-radius: var(--r-sm) !important;
    font-weight: 700 !important;
    font-size: 1rem !important;
    padding: 0.75rem 1.5rem !important;
    box-shadow: 0 4px 18px var(--accent-glow) !important;
    transition: all var(--t) !important;
}
.stDownloadButton > button:hover {
    transform: translateY(-2px) !important;
    box-shadow: 0 8px 24px var(--accent-glow) !important;
}

/* ── Tabs ── */
.stTabs [data-baseweb="tab-list"] {
    background: var(--elevated) !important;
    border-radius: var(--r-sm) !important;
    gap: 2px !important;
    padding: 3px !important;
    border: 1px solid var(--border) !important;
}
.stTabs [data-baseweb="tab"] {
    background: transparent !important;
    color: var(--text2) !important;
    border-radius: 4px !important;
    font-family: inherit !important;
    font-size: 0.82rem !important;
    padding: 0.4rem 0.75rem !important;
}
.stTabs [aria-selected="true"] {
    background: var(--accent) !important;
    color: white !important;
}
.stTabs [data-baseweb="tab-panel"] {
    background: transparent !important;
    padding-top: 1.5rem !important;
}

/* ── Progress ── */
.stProgress > div > div {
    background: var(--accent) !important;
    border-radius: 100px !important;
}
.stProgress > div {
    background: var(--elevated) !important;
    border-radius: 100px !important;
    height: 6px !important;
}

/* ── Misc ── */
hr { border-color: var(--border) !important; margin: 1.5rem 0 !important; }
.stAlert { border-radius: var(--r-sm) !important; font-family: inherit !important; font-size: 0.9rem !important; }

/* 모든 heading 텍스트 색상 (h4, h5, h6 포함) */
.stMarkdown h1, .stMarkdown h2, .stMarkdown h3,
.stMarkdown h4, .stMarkdown h5, .stMarkdown h6 {
    color: var(--text) !important; letter-spacing: -0.3px;
}
.stMarkdown p { color: var(--text2) !important; line-height: 1.65 !important; }

/* caption, 일반 텍스트 */
.stCaption p, [data-testid="stCaptionContainer"] p,
small, .caption { color: var(--text3) !important; }

/* checkbox / radio 라벨 */
.stCheckbox label, .stCheckbox label span,
[data-testid="stCheckbox"] label,
[data-testid="stCheckbox"] p {
    color: var(--text) !important;
}
.stRadio label, .stRadio label span,
[data-testid="stRadio"] label,
[data-testid="stRadio"] p {
    color: var(--text2) !important;
}

/* st.write / 일반 텍스트 */
[data-testid="stText"] { color: var(--text2) !important; }
.stMarkdown { color: var(--text2) !important; }

/* divider */
[data-testid="stDivider"] hr {
    border-color: var(--border) !important;
}

/* ═══════════════════════════════════════
   Custom HTML Components
═══════════════════════════════════════ */

/* ── Navbar ── */
.navbar {
    display: flex; align-items: center; justify-content: space-between;
    padding: 1rem 2rem;
    background: var(--nav-bg);
    backdrop-filter: blur(20px);
    border-bottom: 1px solid var(--border);
}
.nav-logo {
    font-weight: 900; font-size: 1.4rem;
    background: linear-gradient(135deg,#C4B5FD,#A78BFA,#7C3AED);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    letter-spacing: -0.5px;
}
.nav-links { display: flex; gap: 2rem; align-items: center; }
.nav-link { color: var(--text2); text-decoration: none; font-size: 0.87rem; font-weight: 500; transition: color var(--t); }
.nav-link:hover { color: var(--text); }
.nav-cta {
    background: linear-gradient(135deg,#7C3AED,#4F46E5);
    color: white; text-decoration: none;
    font-size: 0.82rem; font-weight: 700;
    padding: 0.45rem 1.1rem; border-radius: 100px;
    box-shadow: 0 4px 14px var(--accent-glow);
    transition: all var(--t);
}
.nav-cta:hover { opacity: 0.88; transform: translateY(-1px); }

/* ── Hero ── */
.hero-outer {
    position: relative; overflow: hidden;
    padding: 7rem 1rem 4.5rem;
}
.hero-glow {
    position: absolute; top: -160px; left: 50%; transform: translateX(-50%);
    width: 1000px; height: 600px;
    background: radial-gradient(ellipse 60% 55% at 50% 0%, rgba(124,58,237,0.14) 0%, transparent 70%);
    pointer-events: none; z-index: 0;
}
.hero-inner { position: relative; z-index: 1; text-align: center; max-width: 820px; margin: 0 auto; }
.hero-badge {
    display: inline-flex; align-items: center; gap: 0.45rem;
    background: rgba(124,58,237,0.1);
    border: 1px solid rgba(124,58,237,0.28); color: var(--accent-l);
    font-size: 0.73rem; font-weight: 700; padding: 0.32rem 1rem;
    border-radius: 100px; margin-bottom: 2rem;
    letter-spacing: 0.08em; text-transform: uppercase;
}
.hero-title {
    font-size: clamp(2.8rem, 5.5vw, 4.4rem);
    font-weight: 900; line-height: 1.08;
    margin-bottom: 1.4rem; color: var(--text);
    letter-spacing: -2.5px;
}
.hero-title .g {
    background: linear-gradient(130deg,#C4B5FD 0%,#A78BFA 35%,#7C3AED 65%,#4F46E5 100%);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
}
.hero-sub {
    font-size: 1.08rem; color: var(--text2); line-height: 1.82;
    margin-bottom: 2.8rem; max-width: 560px; margin-left: auto; margin-right: auto;
}
.hero-proof {
    display: flex; justify-content: center; align-items: center;
    gap: 1.25rem; margin-top: 1.75rem; flex-wrap: wrap;
}
.hero-proof-item {
    display: flex; align-items: center; gap: 0.35rem;
    font-size: 0.82rem; color: var(--text3);
}
.hero-proof-item b { color: var(--text2); font-weight: 700; }
.hero-proof-sep { width: 3px; height: 3px; border-radius: 50%; background: var(--text3); }

/* ── Stats Strip ── */
.stats-strip {
    display: flex; justify-content: center; align-items: center; gap: 0;
    border-top: 1px solid var(--border);
    border-bottom: 1px solid var(--border);
    margin: 0 0 5.5rem;
}
.stat-item {
    flex: 1; text-align: center; padding: 2.2rem 1rem;
    border-right: 1px solid var(--border);
}
.stat-item:last-child { border-right: none; }
.stat-num {
    font-size: 2.1rem; font-weight: 900; letter-spacing: -1.5px;
    background: linear-gradient(135deg,#C4B5FD,#A78BFA,#7C3AED);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    margin-bottom: 4px;
}
.stat-label { font-size: 0.8rem; color: var(--text2); font-weight: 500; }

/* ── How it Works ── */
.steps-grid {
    display: grid; grid-template-columns: repeat(3,1fr); gap: 1.25rem;
    margin-bottom: 5.5rem;
}
.step-card {
    background: var(--surface); border: 1px solid var(--border);
    border-radius: var(--r-lg); padding: 2rem 1.75rem;
    position: relative; overflow: hidden;
    transition: all 0.22s ease;
}
.step-card:hover {
    border-color: rgba(124,58,237,0.32);
    transform: translateY(-4px);
    box-shadow: 0 20px 50px rgba(0,0,0,0.35), 0 0 0 1px rgba(124,58,237,0.1);
}
.step-card::before {
    content: ''; position: absolute; top: 0; left: 0; right: 0; height: 2px;
    background: linear-gradient(90deg, #7C3AED, #4F46E5, #818CF8);
}
.step-num {
    display: inline-flex; align-items: center; justify-content: center;
    width: 42px; height: 42px; border-radius: 50%;
    background: var(--accent-bg); color: var(--accent-l);
    font-size: 0.88rem; font-weight: 900; margin-bottom: 1.25rem;
    border: 1px solid rgba(124,58,237,0.25);
}
.step-title { font-size: 1rem; font-weight: 800; color: var(--text); margin-bottom: 0.6rem; letter-spacing: -0.2px; }
.step-desc { font-size: 0.85rem; color: var(--text2); line-height: 1.75; }

/* ── Features ── */
.slabel {
    display: block; text-align: center; color: var(--accent-l);
    font-size: 0.72rem; font-weight: 700; letter-spacing: 0.12em;
    text-transform: uppercase; margin-bottom: 0.65rem;
}
.stitle {
    text-align: center; font-size: 2rem; font-weight: 900;
    color: var(--text); letter-spacing: -1px; margin-bottom: 0.65rem;
}
.ssub {
    text-align: center; font-size: 0.97rem; color: var(--text2);
    line-height: 1.72; max-width: 560px; margin: 0 auto 2.75rem;
}
.feat-grid {
    display: grid; grid-template-columns: repeat(auto-fit,minmax(240px,1fr));
    gap: 1rem; margin-bottom: 5.5rem;
}
.feat {
    background: var(--feat-bg);
    border: 1px solid var(--border); border-radius: var(--r-lg); padding: 1.75rem;
    transition: all 0.22s ease;
    box-shadow: 0 2px 8px rgba(0,0,0,0.04);
}
.feat:hover {
    border-color: rgba(124,58,237,0.35);
    background: var(--accent-bg);
    transform: translateY(-3px);
    box-shadow: 0 14px 40px rgba(0,0,0,0.1), 0 0 0 1px rgba(124,58,237,0.08);
}
.feat-icon { font-size: 1.6rem; margin-bottom: 0.9rem; display: block; }
.feat-title { font-size: 0.97rem; font-weight: 800; color: var(--text); margin-bottom: 0.45rem; }
.feat-desc { font-size: 0.84rem; color: var(--text2); line-height: 1.68; }

/* ── Mockup ── */
.mockup {
    max-width: 920px; margin: 3rem auto 5.5rem;
    background: var(--surface); border: 1px solid var(--border);
    border-radius: var(--r-lg); overflow: hidden;
    box-shadow:
        0 0 0 1px rgba(255,255,255,0.04),
        0 40px 80px rgba(0,0,0,0.55),
        0 0 80px rgba(124,58,237,0.07);
}
.mockup-bar {
    display: flex; justify-content: space-between; align-items: center;
    padding: 0.7rem 1.25rem; background: var(--elevated); border-bottom: 1px solid var(--border);
}
.mdots { display:flex; gap:5px; }
.md { width:10px; height:10px; border-radius:50%; }
.md-r { background:#EF4444; } .md-y { background:#F59E0B; } .md-g { background:#10B981; }
.mfile { font-size:0.72rem; color:var(--text3); font-family:monospace; }
.msize { font-size:0.78rem; color:var(--accent-l); font-weight:700; }
.mockup-body { background:#0A0F1E; padding:2.5rem; }
.m-slide-hdr { border-left:3px solid #6366F1; padding-left:1.1rem; margin-bottom:1.4rem; }
.m-title { font-size:1.75rem; font-weight:900; color:white; margin-bottom:4px; letter-spacing:-0.5px; }
.m-sub { font-size:0.82rem; color:#94A3B8; }
.m-cards { display:flex; gap:0.875rem; margin-bottom:1.25rem; }
.m-card {
    flex:1; background:rgba(30,41,59,0.8);
    border:1px solid rgba(51,65,85,0.7); border-radius:8px; padding:0.9rem;
    backdrop-filter: blur(4px);
}
.m-ctitle { font-size:0.78rem; font-weight:700; color:#818CF8; margin-bottom:5px; }
.m-cdesc { font-size:0.72rem; color:#CBD5E1; line-height:1.5; }
.m-tw { background:rgba(30,27,75,0.6); border:1px solid rgba(49,46,129,0.6); border-radius:6px; padding:0.6rem 0.9rem; font-size:0.8rem; color:#A5B4FC; }

/* ── Pricing ── */
.price-wrap { display:flex; gap:1.25rem; max-width:680px; margin:0 auto 5.5rem; }
.plan {
    flex:1; background:var(--surface); border:1px solid var(--border);
    border-radius:var(--r-lg); padding:2rem;
    transition: all 0.22s ease;
}
.plan:hover { border-color: var(--border-b); transform: translateY(-2px); }
.plan.pro {
    border-color:rgba(124,58,237,0.45);
    background:linear-gradient(160deg,rgba(124,58,237,0.09),rgba(79,70,229,0.04));
    position:relative;
}
.plan.pro:hover { border-color: rgba(124,58,237,0.65); }
.plan-badge {
    position:absolute; top:-12px; left:50%; transform:translateX(-50%);
    background:linear-gradient(135deg,#7C3AED,#4F46E5);
    color:white; font-size:0.63rem; font-weight:700; padding:4px 14px;
    border-radius:100px; letter-spacing:0.08em; text-transform:uppercase; white-space:nowrap;
    box-shadow: 0 4px 12px rgba(124,58,237,0.4);
}
.plan-name { font-size:0.9rem; font-weight:600; color:var(--text2); margin-bottom:0.45rem; }
.plan-price { font-size:2.6rem; font-weight:900; color:var(--text); line-height:1; margin-bottom:0.2rem; letter-spacing:-1.5px; }
.plan-unit { font-size:0.85rem; font-weight:400; color:var(--text3); }
.plan-list { list-style:none; padding:0; margin:1.5rem 0 0; }
.plan-list li {
    font-size:0.85rem; color:var(--text2); padding:0.38rem 0;
    display:flex; gap:0.5rem; align-items:flex-start;
    border-bottom:1px solid var(--border);
}
.plan-list li:last-child { border-bottom:none; }
.plan-list li::before { content:'✓'; color:var(--green); font-weight:700; flex-shrink:0; margin-top:1px; }

/* ── CTA Banner ── */
.cta-banner {
    position: relative; overflow: hidden;
    background: linear-gradient(145deg, rgba(124,58,237,0.13), rgba(79,70,229,0.06));
    border: 1px solid rgba(124,58,237,0.22);
    border-radius: var(--r-lg); padding: 4.5rem 2rem;
    text-align: center; margin-bottom: 4rem;
}
.cta-banner::before {
    content: ''; position: absolute; top: -140px; left: 50%; transform: translateX(-50%);
    width: 700px; height: 450px;
    background: radial-gradient(ellipse, rgba(124,58,237,0.18), transparent 70%);
    pointer-events: none;
}
.cta-banner-inner { position: relative; z-index: 1; }
.cta-banner-title {
    font-size: 2.1rem; font-weight: 900; color: var(--text);
    letter-spacing: -1px; margin-bottom: 0.8rem;
}
.cta-banner-sub { font-size: 1rem; color: var(--text2); margin-bottom: 2.25rem; line-height: 1.7; }

/* ── Auth ── */
.auth-box { max-width:400px; margin:3rem auto; background:var(--surface); border:1px solid var(--border); border-radius:var(--r-lg); padding:2.25rem; box-shadow:0 24px 48px rgba(0,0,0,0.4); }
.auth-title { text-align:center; font-size:1.55rem; font-weight:900; color:var(--text); margin-bottom:0.35rem; letter-spacing:-0.5px; }
.auth-sub { text-align:center; font-size:0.86rem; color:var(--text2); margin-bottom:1.75rem; line-height:1.6; }
.divline { display:flex; align-items:center; gap:0.65rem; color:var(--text3); font-size:0.78rem; margin:0.9rem 0; }
.divline::before,.divline::after { content:''; flex:1; height:1px; background:var(--border); }

/* ── Wizard Progress Bar ── */
.wiz-bar { display:flex; align-items:center; padding:0.875rem 1.5rem; background:var(--surface); border:1px solid var(--border); border-radius:var(--r-md); margin-bottom:1.75rem; }
.w-step { display:flex; align-items:center; gap:0.45rem; flex-shrink:0; }
.w-num { width:24px; height:24px; border-radius:50%; display:flex; align-items:center; justify-content:center; font-size:0.7rem; font-weight:700; border:2px solid var(--border); color:var(--text3); background:var(--bg); transition:all var(--t); }
.w-label { font-size:0.78rem; font-weight:500; color:var(--text3); white-space:nowrap; }
.w-step.active .w-num { background:var(--accent); border-color:var(--accent); color:white; }
.w-step.active .w-label { color:var(--text); font-weight:600; }
.w-step.done .w-num { background:var(--green-bg); border-color:var(--green); color:var(--green); }
.w-step.done .w-label { color:var(--green); }
.w-line { flex:1; height:1px; background:var(--border); margin:0 0.4rem; min-width:10px; }
.w-line.done { background:var(--green); }

/* ── Slide Edit Cards (Step 3) ── */
.s-card { background:var(--surface); border:1px solid var(--border); border-left:3px solid var(--accent); border-radius:var(--r-md); padding:1.4rem 1.5rem; margin-bottom:1.1rem; }
.s-card-hdr { display:flex; justify-content:space-between; align-items:center; margin-bottom:1rem; }
.s-badge { font-size:0.7rem; font-weight:700; color:var(--accent-l); text-transform:uppercase; letter-spacing:0.07em; }
.s-type { font-size:0.68rem; background:var(--accent-bg); color:var(--accent-l); border:1px solid rgba(124,58,237,0.2); padding:2px 7px; border-radius:4px; }

/* ── Canvas Preview (Step 5) ── */
.canvas { background:#0A0F1E; border:2px solid #1E293B; border-radius:var(--r-md); padding:2rem 2.25rem; position:relative; min-height:320px; }
.canvas-lbl { position:absolute; top:10px; right:12px; background:#1E293B; color:#818CF8; border:1px solid #334155; padding:3px 8px; border-radius:4px; font-size:0.68rem; font-family:monospace; }
.c-hdr { border-left:3px solid #6366F1; padding-left:1rem; margin:1.25rem 0; }
.c-t { font-size:1.65rem; font-weight:900; color:white; line-height:1.2; margin-bottom:0.3rem; letter-spacing:-0.5px; }
.c-s { font-size:0.83rem; color:#94A3B8; font-style:italic; }
.c-tw { background:rgba(30,27,75,0.6); border:1px solid rgba(49,46,129,0.6); border-radius:6px; padding:0.6rem 0.9rem; color:#A5B4FC; font-size:0.85rem; margin-top:0.875rem; }
.c-b { font-size:0.95rem; color:#E2E8F0; margin-bottom:9px; }
.c-col { flex:1; background:#1E293B; border:1px solid #334155; border-radius:6px; padding:0.8rem; }
.c-col-lbl { font-size:0.72rem; font-weight:700; color:#818CF8; margin-bottom:5px; }

/* ── Slide Miniature (Step 4) ── */
.mini { background:white; border:1px solid #E2E8F0; border-radius:8px; padding:10px; min-height:120px; transition:border-color var(--t); box-shadow:0 2px 6px rgba(0,0,0,0.07); }
.mini:hover { border-color:#A78BFA; box-shadow: 0 4px 14px rgba(124,58,237,0.12); }
.mini-n { font-size:0.62rem; font-weight:700; color:#6E7681; font-family:monospace; margin-bottom:3px; }
.mini-t { font-size:0.78rem; font-weight:700; color:#1E293B; line-height:1.3; margin-bottom:3px; }
.mini-tp { font-size:0.62rem; color:#7C3AED; margin-bottom:5px; }
.mini-bar { width:22px; height:2px; background:#7C3AED; margin-bottom:5px; }
.mini-kw { font-size:0.6rem; color:#4F46E5; font-style:italic; line-height:1.4; }

/* ── Info tag ── */
.tag-info { display:inline-block; background:rgba(88,166,255,0.08); border:1px solid rgba(88,166,255,0.2); color:#58A6FF; font-size:0.82rem; padding:0.4rem 0.8rem; border-radius:var(--r-sm); margin-bottom:1rem; }

/* ── Agent Pipeline ── */
.agent-pipe { background:var(--surface); border:1px solid var(--border); border-radius:var(--r-lg); padding:1.5rem 2rem; margin-bottom:1.75rem; }
.agent-pipe-title { font-size:0.7rem; font-weight:700; color:var(--accent-l); text-transform:uppercase; letter-spacing:0.12em; margin-bottom:1.25rem; }
.agent-row { display:flex; align-items:center; gap:1rem; padding:0.65rem 0; border-bottom:1px solid var(--border); }
.agent-row:last-child { border-bottom:none; }
.agent-dot { width:10px; height:10px; border-radius:50%; background:var(--text3); flex-shrink:0; transition:all 0.3s; }
.agent-dot.done { background:var(--green); }
.agent-dot.active { background:var(--accent); box-shadow:0 0 10px var(--accent-glow); animation: pulse 1.2s infinite; }
@keyframes pulse { 0%,100%{opacity:1;} 50%{opacity:0.5;} }
.agent-label { font-size:0.88rem; font-weight:700; color:var(--text); flex:1; }
.agent-sub { font-size:0.75rem; color:var(--text3); }
.agent-status { font-size:0.75rem; font-weight:600; color:var(--text3); white-space:nowrap; }
.agent-status.done { color:var(--green); }
.agent-status.active { color:var(--accent-l); }

/* ── Structure Cards (Stage 1 결과) ── */
.struct-grid { display:grid; grid-template-columns:repeat(auto-fill,minmax(260px,1fr)); gap:0.875rem; margin:1.5rem 0 2rem; }
.struct-card { background:var(--surface); border:1px solid var(--border); border-radius:var(--r-md); padding:1.1rem 1.2rem; display:flex; gap:0.875rem; align-items:flex-start; transition:all 0.2s; }
.struct-card:hover { border-color:rgba(124,58,237,0.3); box-shadow:0 4px 14px rgba(0,0,0,0.08); }
.struct-num { width:32px; height:32px; border-radius:50%; background:var(--accent-bg); color:var(--accent-l); display:flex; align-items:center; justify-content:center; font-size:0.72rem; font-weight:800; flex-shrink:0; border:1px solid rgba(124,58,237,0.2); }
.struct-body { flex:1; min-width:0; }
.struct-title { font-size:0.88rem; font-weight:700; color:var(--text); margin-bottom:4px; line-height:1.35; }
.struct-badge { font-size:0.65rem; background:var(--accent-bg); color:var(--accent-l); padding:2px 7px; border-radius:4px; display:inline-block; margin-bottom:5px; }
.struct-role { font-size:0.77rem; color:var(--text2); line-height:1.5; }

/* ── Step 2: Format & Theme Cards ── */
.fmt-card {
    border: 2px solid var(--border);
    border-radius: 12px;
    padding: 10px 8px 6px;
    text-align: center;
    background: var(--surface);
    transition: border-color 0.18s, box-shadow 0.18s;
    cursor: pointer;
    margin-bottom: 6px;
}
.fmt-card.sel { border-color: var(--accent); box-shadow: 0 0 0 3px var(--accent-glow); }
.fmt-ratio-box {
    width: 100%;
    background: var(--elevated);
    border-radius: 6px;
    margin-bottom: 7px;
    position: relative;
    overflow: hidden;
}
.fmt-ratio-box::before { content:''; display:block; }
.fmt-ratio-inner { position:absolute; inset:0; display:flex; align-items:center; justify-content:center; }
.fmt-ratio-inner span { font-size:0.68rem; font-weight:700; color:var(--text2); }
.fmt-label { font-size:0.78rem; font-weight:600; color:var(--text); }
.fmt-desc  { font-size:0.68rem; color:var(--text2); margin-top:2px; }

.theme-card {
    border: 2px solid var(--border);
    border-radius: 12px;
    overflow: hidden;
    background: var(--surface);
    transition: border-color 0.18s, box-shadow 0.18s;
    cursor: pointer;
    margin-bottom: 6px;
}
.theme-card.sel { border-color: var(--accent); box-shadow: 0 0 0 3px var(--accent-glow); }
.theme-preview { height: 72px; padding: 10px 12px; position: relative; overflow: hidden; }
.theme-bar { height:4px; border-radius:2px; margin-bottom:7px; width:55%; }
.theme-line { height:2px; border-radius:1px; margin-bottom:4px; }
.theme-dot { position:absolute; bottom:8px; right:8px; width:18px; height:18px; border-radius:50%; opacity:0.85; }
.theme-foot { padding:5px 10px 6px; display:flex; align-items:center; justify-content:space-between; background:var(--surface); }
.theme-name { font-size:0.76rem; font-weight:600; color:var(--text); }
.theme-tag  { font-size:0.64rem; background:var(--accent-bg); color:var(--accent); padding:1px 6px; border-radius:10px; }
</style>
""", unsafe_allow_html=True)

# ── Session State ─────────────────────────────────────────────────────────────
_DEFAULTS = {
    'view': 'landing',
    'step': 1,
    'topic': '',
    'source_text': '',
    'lang': '한국어',
    'category': 'Startup Pitch (투자 유치 및 IR 덱)',
    'deck_format': 'Standard (16:9)',
    'template_style': 'Modern',
    'theme': 'Nordic Tech (Ice White & Indigo)',
    'presentation_data': None,
    'slide_count': 7,
    'generation_progress': False,
    'compiled_pptx': None,
    'chat_history': [],
    'gemini_api_key': '',
    'dark_mode': False,
    'agent_stage': 0,
    'structure_data': None,
    'custom_accent': '#7C3AED',
    'custom_bg': '#F8F9FF',
    'custom_text': '#0F0F2A',
}
for _k, _v in _DEFAULTS.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v

# 다크 모드 CSS 주입
if st.session_state.get('dark_mode', False):
    st.markdown("""<style>
    :root {
        --bg: #080810; --surface: #10101C; --elevated: #181828;
        --border: rgba(255,255,255,0.07); --border-b: rgba(255,255,255,0.14);
        --text: #EEEEFF; --text2: #8888AA; --text3: #44445A;
        --accent-l: #A78BFA;
        --accent-bg: rgba(124,58,237,0.12); --accent-glow: rgba(124,58,237,0.28);
        --green: #22C55E; --green-bg: rgba(34,197,94,0.1);
        --nav-bg: rgba(8,8,16,0.92); --feat-bg: rgba(255,255,255,0.03);
    }
    </style>""", unsafe_allow_html=True)

# ── API ───────────────────────────────────────────────────────────────────────
def initialize_gemini():
    key = os.getenv("GEMINI_API_KEY") or st.session_state.get('gemini_api_key', '')
    if key:
        genai.configure(api_key=key)
    return key

def get_groq_client():
    key = os.getenv("GROQ_API_KEY") or st.session_state.get('groq_api_key', '')
    if not key:
        return None
    return Groq(api_key=key)

api_key = initialize_gemini()
groq_client = get_groq_client()


def ask_gemini_for_presentation(topic, lang, category, slide_count, source_text=''):
    if not initialize_gemini():
        st.error("Gemini API Key가 필요합니다.")
        return None
    model = genai.GenerativeModel("gemini-1.5-flash")
    src = f"\n\n[참고 자료 (요약해서 활용하세요)]\n{source_text[:3500]}" if source_text.strip() else ''
    prompt = f"""
주제: "{topic}"
언어: {lang}
발표 목적: {category}
슬라이드 수: {slide_count}장
{src}

위 정보를 바탕으로 {slide_count}장 분량의 고품질 발표 슬라이드를 {lang}로 기획해주세요.

사용 가능한 슬라이드 타입:
- title_and_content : 핵심 불렛 슬라이드 (bullets 3~4개)
- section_header    : 섹션 전환 슬라이드
- big_stat          : 수치/지표 강조 슬라이드
- two_column        : 비교/대조 슬라이드
- three_cards       : 3개 카드 슬라이드
- timeline          : 로드맵/단계 슬라이드 (4단계)
- team_grid         : 팀 소개 슬라이드

엄격히 지킬 규칙:
1. 모든 텍스트는 {lang}로 작성
2. big_stat  → stat_value (예: "85%"), stat_description 필수 포함
3. three_cards → cards 배열에 card_title, card_content 3개 필수
4. timeline   → timeline_steps 배열에 step_title, step_desc 4개 필수
5. team_grid  → team_members 배열에 name, role, bio 필수
6. bullets    → 3~4개, 간결한 발표용 문장
7. key_takeaway → 슬라이드 핵심 인사이트 1문장
8. speaker_notes → 발표자용 스크립트 2~3문장 (청중에게 말할 내용)

순수 JSON만 반환 (```코드블록 없이):
{{
  "presentation_title": "...",
  "presentation_subtitle": "...",
  "slides": [
    {{
      "slide_index": 1,
      "slide_type": "title_and_content",
      "title": "...",
      "summary": "1문장 요약",
      "bullets": ["..."],
      "key_takeaway": "...",
      "speaker_notes": "발표자 스크립트..."
    }}
  ]
}}
"""
    cfg = {"response_mime_type": "application/json", "temperature": 0.75}
    try:
        resp = model.generate_content(prompt, generation_config=cfg)
        text = resp.text.strip()
        text = re.sub(r'^```(?:json)?\n?', '', text)
        text = re.sub(r'\n?```$', '', text)
        data = json.loads(text)
        # Assign unique IDs to each slide for stable widget keys
        for s in data.get('slides', []):
            s['_id'] = uuid.uuid4().hex[:8]
        return data
    except Exception as e:
        st.error(f"AI 생성 오류: {e}")
        return None


# ── Agent 1: 구조 설계 (Groq — 빠름) ────────────────────────────
def run_structure_agent(topic, lang, category, slide_count, source_text=''):
    src = f"\n[참고 자료 요약]\n{source_text[:2000]}" if source_text.strip() else ''
    prompt = f"""당신은 프레젠테이션 구조 전문가입니다.
아래 발표의 최적 슬라이드 구조만 빠르게 설계하세요. 상세 내용은 불필요합니다.

주제: {topic}
발표 목적: {category}
언어: {lang}
슬라이드 수: {slide_count}장
{src}

요구사항:
- 스토리 흐름이 자연스럽도록 구성 (문제→해결→증거→행동촉구 등)
- 각 슬라이드에 가장 적합한 타입 선택
- 제목은 {lang}로 작성

타입: title_and_content, section_header, big_stat, two_column, three_cards, timeline, team_grid

순수 JSON만 반환 (코드블록 없이):
{{"presentation_title":"...","presentation_subtitle":"...","slides":[{{"slide_index":1,"slide_type":"title_and_content","title":"...","role":"이 슬라이드가 담당하는 역할 1문장"}}]}}"""

    client = get_groq_client()
    if client:
        try:
            resp = client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.6,
                response_format={"type": "json_object"},
                max_tokens=3000,
            )
            data = json.loads(resp.choices[0].message.content)
            for s in data.get('slides', []):
                s['_id'] = uuid.uuid4().hex[:8]
            return data
        except Exception as e:
            st.warning(f"Groq 오류 ({e}), Gemini로 대체합니다.")

    # Gemini 폴백
    if not initialize_gemini():
        st.error("API Key가 없습니다.")
        return None
    model = genai.GenerativeModel("gemini-1.5-flash")
    cfg = {"response_mime_type": "application/json", "temperature": 0.6}
    try:
        resp = model.generate_content(prompt, generation_config=cfg)
        text = re.sub(r'^```(?:json)?\n?', '', resp.text.strip())
        text = re.sub(r'\n?```$', '', text)
        data = json.loads(text)
        for s in data.get('slides', []):
            s['_id'] = uuid.uuid4().hex[:8]
        return data
    except Exception as e:
        st.error(f"[구조 설계 오류] {e}")
        return None


# ── Agent 2: 콘텐츠 생성 (Groq 우선) ────────────────────────────
def run_content_agent(structure, topic, lang, category, source_text=''):
    src = f"\n[참고 자료]\n{source_text[:3000]}" if source_text.strip() else ''
    structure_str = json.dumps(
        [{"slide_index": s.get("slide_index", i+1), "slide_type": s.get("slide_type", "title_and_content"),
          "title": s.get("title", ""), "_id": s.get("_id", "")}
         for i, s in enumerate(structure.get("slides", []))],
        ensure_ascii=False, indent=2
    )
    prompt = f"""
당신은 프레젠테이션 콘텐츠 전문가입니다.
아래 확정된 구조에 맞는 {lang} 콘텐츠를 생성하세요.

[슬라이드 구조]
{structure_str}

[발표 맥락]
주제: {topic} / 목적: {category}
{src}

각 슬라이드 타입별 필수 필드:
- title_and_content / two_column: bullets 3~4개
- section_header: summary 1문장
- big_stat: stat_value("85%" 형식), stat_description, bullets 2~3개
- three_cards: cards [{{"card_title":"...", "card_content":"..."}} × 3]
- timeline: timeline_steps [{{"step_title":"...", "step_desc":"..."}} × 4]
- team_grid: team_members [{{"name":"...", "role":"...", "bio":"..."}} × 2~3]
- 모든 슬라이드: key_takeaway(핵심 1문장), speaker_notes(발표 스크립트 2~3문장)
- _id 값은 구조에서 그대로 유지

순수 JSON만 반환:
{{
  "presentation_title": "{structure.get('presentation_title', '')}",
  "presentation_subtitle": "{structure.get('presentation_subtitle', '')}",
  "slides": [...]
}}
"""
    def _restore_ids(data):
        """Structure agent가 생성한 _id를 slide_index 기준으로 복원."""
        id_map = {s.get('slide_index'): s.get('_id', '') for s in structure.get('slides', [])}
        for s in data.get('slides', []):
            if not s.get('_id'):
                s['_id'] = id_map.get(s.get('slide_index')) or uuid.uuid4().hex[:8]
        return data

    client = get_groq_client()
    if client:
        try:
            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.75,
                response_format={"type": "json_object"},
                max_tokens=6000,
            )
            return _restore_ids(json.loads(resp.choices[0].message.content))
        except Exception as e:
            st.warning(f"Groq 오류 ({e}), Gemini로 대체합니다.")

    if not initialize_gemini():
        st.error("API Key가 없습니다. 설정에서 Gemini API Key를 입력해 주세요.")
        return None
    model = genai.GenerativeModel("gemini-1.5-flash")
    cfg = {"response_mime_type": "application/json", "temperature": 0.75}
    try:
        resp = model.generate_content(prompt, generation_config=cfg)
        text = re.sub(r'^```(?:json)?\n?', '', resp.text.strip())
        text = re.sub(r'\n?```$', '', text)
        return _restore_ids(json.loads(text))
    except Exception as e:
        st.error(f"[콘텐츠 생성 오류] {e}")
        return None


# ── Agent 3: 품질 개선 (Groq 우선) ──────────────────────────────
def run_polish_agent(data, lang):
    slides_str = json.dumps(data.get('slides', []), ensure_ascii=False, indent=2)[:6000]
    prompt = f"""
당신은 프레젠테이션 품질 전문가입니다.
다음 슬라이드 덱을 검토하고 텍스트 품질만 개선하세요.

[슬라이드]
{slides_str}

개선 지침:
1. key_takeaway: 더 강렬하고 기억에 남는 1문장으로
2. speaker_notes: 자연스러운 {lang} 발표 구어체로
3. bullets: 병렬 구조 통일, 각 항목 간결하게
4. 전체 톤을 {lang} 전문 발표에 맞게 일관되게
5. 절대 변경 금지: slide_type, slide_index, _id, stat_value, cards 구조, timeline_steps 구조

개선된 slides 배열만 JSON 반환:
{{"slides": [...]}}
"""
    def _apply_polish(result):
        if 'slides' in result:
            for s in result['slides']:
                if '_id' not in s:
                    orig = next((x for x in data['slides'] if x.get('slide_index') == s.get('slide_index')), None)
                    if orig:
                        s['_id'] = orig.get('_id', uuid.uuid4().hex[:8])
            data['slides'] = result['slides']
        return data

    client = get_groq_client()
    if client:
        try:
            resp = client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.45,
                response_format={"type": "json_object"},
                max_tokens=6000,
            )
            return _apply_polish(json.loads(resp.choices[0].message.content))
        except Exception as e:
            st.warning(f"Groq 오류 ({e}), Gemini로 대체합니다.")

    if not initialize_gemini():
        return data
    model = genai.GenerativeModel("gemini-1.5-flash")
    cfg = {"response_mime_type": "application/json", "temperature": 0.45}
    try:
        resp = model.generate_content(prompt, generation_config=cfg)
        text = re.sub(r'^```(?:json)?\n?', '', resp.text.strip())
        text = re.sub(r'\n?```$', '', text)
        return _apply_polish(json.loads(text))
    except Exception as e:
        st.error(f"[품질 개선 오류] {e}")
        return data


def magic_edit_slide(slide, instruction, lang):
    prompt = f"""현재 슬라이드 데이터:
{json.dumps(slide, ensure_ascii=False, indent=2)}

수정 지시사항: {instruction}
언어: {lang}

지시사항에 맞게 슬라이드를 수정하여 동일한 JSON 구조로 반환하세요.
slide_type, slide_index, _id 는 변경하지 마세요.
순수 JSON만 반환 (코드블록 없이)."""

    client = get_groq_client()
    if client:
        try:
            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.5,
                response_format={"type": "json_object"},
                max_tokens=3000,
            )
            return json.loads(resp.choices[0].message.content)
        except Exception as e:
            st.warning(f"Groq 오류 ({e}), Gemini로 대체합니다.")

    if not initialize_gemini():
        st.error("Magic Edit에 API Key가 필요합니다.")
        return None
    model = genai.GenerativeModel("gemini-1.5-flash")
    cfg = {"response_mime_type": "application/json", "temperature": 0.5}
    try:
        resp = model.generate_content(prompt, generation_config=cfg)
        text = re.sub(r'^```(?:json)?\n?', '', resp.text.strip())
        text = re.sub(r'\n?```$', '', text)
        return json.loads(text)
    except Exception as e:
        st.error(f"Magic Edit 오류: {e}")
        return None


def _hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _get_colors(theme_key):
    if theme_key == "Custom":
        return (st.session_state.get('custom_accent', '#7C3AED'),
                st.session_state.get('custom_bg', '#F8F9FF'),
                st.session_state.get('custom_text', '#0F0F2A'))
    c = THEMES.get(theme_key, list(THEMES.values())[0])
    return ("#{:02X}{:02X}{:02X}".format(*c["accent"]),
            "#{:02X}{:02X}{:02X}".format(*c["bg_light"]),
            "#{:02X}{:02X}{:02X}".format(*c["primary_text"]))


def _get_dark_bg_hex(theme_key):
    if theme_key == "Custom":
        return st.session_state.get('custom_bg', '#0F172A')
    c = THEMES.get(theme_key, list(THEMES.values())[0])
    return "#{:02X}{:02X}{:02X}".format(*c["bg_dark"])


def _slide_html(acc, bg, txt, w=16, h=9, radius=8, shadow=True):
    pt = h / w * 100
    sh = "0 4px 20px rgba(0,0,0,0.18)" if shadow else "none"
    return f"""
<div style="width:100%;padding-top:{pt:.2f}%;background:{bg};border-radius:{radius}px;
            position:relative;overflow:hidden;box-shadow:{sh};
            border:1px solid rgba(128,128,128,0.12);">
  <div style="position:absolute;inset:0;">
    <div style="background:{acc};height:20%;display:flex;align-items:center;
                padding:0 7%;gap:6%;">
      <div style="background:rgba(255,255,255,0.92);height:35%;border-radius:2px;width:40%;"></div>
      <div style="flex:1;display:flex;justify-content:flex-end;gap:3%;padding-right:1%;">
        <div style="width:8%;padding-top:8%;border-radius:50%;background:rgba(255,255,255,0.45);"></div>
        <div style="width:8%;padding-top:8%;border-radius:50%;background:rgba(255,255,255,0.25);"></div>
      </div>
    </div>
    <div style="padding:4% 7% 0;">
      <div style="background:{txt};opacity:0.88;height:4px;border-radius:2px;width:58%;margin-bottom:5%;"></div>
      <div style="display:flex;align-items:center;gap:3%;margin-bottom:4%;">
        <div style="width:4%;padding-top:4%;border-radius:50%;background:{acc};flex-shrink:0;opacity:0.9;"></div>
        <div style="background:{txt};opacity:0.32;height:3px;border-radius:1px;flex:1;"></div>
      </div>
      <div style="display:flex;align-items:center;gap:3%;margin-bottom:4%;">
        <div style="width:4%;padding-top:4%;border-radius:50%;background:{acc};flex-shrink:0;opacity:0.9;"></div>
        <div style="background:{txt};opacity:0.22;height:3px;border-radius:1px;width:80%;flex:1;"></div>
      </div>
      <div style="display:flex;align-items:center;gap:3%;">
        <div style="width:4%;padding-top:4%;border-radius:50%;background:{acc};flex-shrink:0;opacity:0.9;"></div>
        <div style="background:{txt};opacity:0.16;height:3px;border-radius:1px;width:88%;flex:1;"></div>
      </div>
    </div>
    <div style="position:absolute;bottom:6%;right:5%;width:12%;padding-top:12%;
                border-radius:50%;background:{acc};opacity:0.15;"></div>
    <div style="position:absolute;bottom:6%;right:18%;width:6%;padding-top:6%;
                border-radius:3px;background:{acc};opacity:0.25;"></div>
  </div>
</div>"""


def _render_slide_pptx_chrome(slide_html, slide_num, total):
    return (
        f'<div style="background:linear-gradient(135deg,#383838,#2a2a2a);border-radius:12px;'
        f'padding:18px 22px 12px;box-shadow:0 8px 28px rgba(0,0,0,0.28);">'
        f'{slide_html}'
        f'<div style="text-align:right;font-size:0.62rem;color:#888;margin-top:8px;'
        f'font-family:monospace;letter-spacing:0.06em;">SLIDE {slide_num} / {total}</div>'
        f'</div>'
    )


def _render_slide_preview_html(slide, acc='#4F46E5', bg='#F8F9FF', txt='#1e293b', style='Modern', dark_bg='#0F172A', large=False):
    def esc(s):
        return str(s or '').replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

    m = 2.0 if large else 1.0  # font scale multiplier for full-width display

    stype  = slide.get('slide_type', 'title_and_content')
    title  = esc(slide.get('title', '') or '')

    # Per-style rendering config
    if style == 'Bold':
        header_h = '25%'
        header_bar = (
            f'<div style="background:{acc};height:{header_h};display:flex;align-items:center;'
            f'padding:0 6%;position:absolute;top:0;left:0;right:0;">'
            f'<div style="font-size:{0.82*m:.2f}rem;font-weight:800;color:white;'
            f'overflow:hidden;white-space:nowrap;text-overflow:ellipsis;max-width:100%;">{title}</div></div>'
        )
        actual_bg     = bg
        content_top   = '27%'
        bullet_dot    = f'<span style="color:{acc};font-weight:700;margin-right:3px;">▶</span>'
        txt_color     = txt
    elif style == 'Minimal':
        header_h = '0%'
        header_bar = (
            f'<div style="position:absolute;top:0;left:0;right:0;padding:5% 6% 0;">'
            f'<div style="font-size:{0.82*m:.2f}rem;font-weight:800;color:{acc};margin-bottom:3px;">{title}</div>'
            f'<div style="height:2px;background:{acc};border-radius:1px;"></div></div>'
        )
        actual_bg     = bg
        content_top   = '22%'
        bullet_dot    = f'<span style="color:{txt};opacity:0.5;margin-right:3px;">—</span>'
        txt_color     = txt
    elif style == 'Dark':
        header_h = '22%'
        header_bar = (
            f'<div style="background:{acc};height:{header_h};display:flex;align-items:center;'
            f'padding:0 6%;position:absolute;top:0;left:0;right:0;">'
            f'<div style="font-size:{0.75*m:.2f}rem;font-weight:700;color:white;'
            f'overflow:hidden;white-space:nowrap;text-overflow:ellipsis;max-width:100%;">{title}</div></div>'
        )
        actual_bg     = dark_bg
        content_top   = '24%'
        bullet_dot    = f'<span style="color:{acc};margin-right:3px;">◆</span>'
        txt_color     = '#e2e8f0'
    else:  # Modern
        header_h = '18%'
        header_bar = (
            f'<div style="background:{acc};height:{header_h};display:flex;align-items:center;'
            f'padding:0 6%;position:absolute;top:0;left:0;right:0;">'
            f'<div style="font-size:{0.75*m:.2f}rem;font-weight:700;color:white;'
            f'overflow:hidden;white-space:nowrap;text-overflow:ellipsis;max-width:100%;">{title}</div></div>'
        )
        actual_bg     = bg
        content_top   = '20%'
        bullet_dot    = f'<div style="width:{5*m:.0f}px;height:{5*m:.0f}px;border-radius:50%;background:{acc};margin-top:3px;flex-shrink:0;"></div>'
        txt_color     = txt

    if stype == 'section_header':
        summary = esc(slide.get('summary', '') or '')
        if style == 'Minimal':
            sec_bg = bg; sec_title_c = acc; sec_sub_c = txt
        elif style == 'Dark':
            sec_bg = dark_bg; sec_title_c = acc; sec_sub_c = '#94a3b8'
        else:
            sec_bg = acc; sec_title_c = 'white'; sec_sub_c = 'rgba(255,255,255,0.75)'
        inner = (
            f'<div style="position:absolute;inset:0;background:{sec_bg};display:flex;flex-direction:column;'
            f'justify-content:center;align-items:center;text-align:center;padding:6%;">'
            f'<div style="font-size:{1.0*m:.2f}rem;font-weight:800;color:{sec_title_c};line-height:1.3;margin-bottom:6px;">{title}</div>'
            f'<div style="width:40px;height:3px;background:{acc};border-radius:2px;margin:8px 0;"></div>'
            f'<div style="font-size:{0.6*m:.2f}rem;color:{sec_sub_c};max-width:80%;">{summary}</div></div>'
        )
    elif stype == 'big_stat':
        sv = esc(slide.get('stat_value', '') or '')
        sd = esc(slide.get('stat_description', '') or '')
        bullets = slide.get('bullets', [])[:2]
        b_html = ''.join([f'<div style="font-size:{0.52*m:.2f}rem;color:{txt_color};opacity:0.7;margin-bottom:2px;">{esc(b)}</div>' for b in bullets])
        inner = (
            f'{header_bar}'
            f'<div style="position:absolute;top:{content_top};left:0;right:0;bottom:0;padding:4% 7%;">'
            f'<div style="font-size:{2.0*m:.2f}rem;font-weight:900;color:{acc};line-height:1;">{sv}</div>'
            f'<div style="font-size:{0.62*m:.2f}rem;color:{txt_color};opacity:0.75;margin:3px 0 7px;">{sd}</div>'
            f'{b_html}</div>'
        )
    elif stype == 'three_cards':
        cards = slide.get('cards', [])[:3]
        if not cards:
            cards = [{"card_title": f"Point {k+1}", "card_content": ""} for k in range(3)]
        card_bg = f'{acc}1A' if style != 'Dark' else f'{acc}33'
        ch = ''.join([
            f'<div style="flex:1;background:{card_bg};border-radius:5px;padding:5px 4px;">'
            f'<div style="font-size:{0.56*m:.2f}rem;font-weight:700;color:{acc};margin-bottom:2px;">{esc(c.get("card_title",""))}</div>'
            f'<div style="font-size:{0.5*m:.2f}rem;color:{txt_color};opacity:0.85;">{esc((c.get("card_content","") or "")[:40])}</div></div>'
            for c in cards])
        inner = f'{header_bar}<div style="position:absolute;top:{content_top};left:0;right:0;bottom:0;padding:4% 5%;display:flex;gap:4px;">{ch}</div>'
    elif stype == 'timeline':
        steps = slide.get('timeline_steps', [])[:4]
        if not steps:
            steps = [{"step_title": f"Phase {k+1}"} for k in range(4)]
        dots = ''.join([
            f'<div style="flex:1;text-align:center;position:relative;z-index:1;">'
            f'<div style="width:{14*m:.0f}px;height:{14*m:.0f}px;border-radius:50%;background:{acc};margin:0 auto 4px;'
            f'display:flex;align-items:center;justify-content:center;">'
            f'<span style="font-size:{0.4*m:.2f}rem;color:white;font-weight:700;">{j+1}</span></div>'
            f'<div style="font-size:{0.5*m:.2f}rem;font-weight:600;color:{txt_color};line-height:1.3;">{esc(s.get("step_title",""))}</div></div>'
            for j, s in enumerate(steps)])
        inner = (
            f'{header_bar}'
            f'<div style="position:absolute;top:{content_top};left:0;right:0;bottom:0;display:flex;align-items:center;padding:0 5%;">'
            f'<div style="width:100%;position:relative;">'
            f'<div style="position:absolute;top:{7*m:.0f}px;left:7%;right:7%;height:1.5px;background:{acc};opacity:0.35;"></div>'
            f'<div style="display:flex;">{dots}</div></div></div>'
        )
    elif stype == 'team_grid':
        members = slide.get('team_members', [])[:3]
        if not members:
            members = [{"name": "Member", "role": ""} for _ in range(2)]
        av = int(22 * m)
        m_html = ''.join([
            f'<div style="flex:1;text-align:center;">'
            f'<div style="width:{av}px;height:{av}px;border-radius:50%;background:{acc};margin:0 auto 3px;'
            f'display:flex;align-items:center;justify-content:center;">'
            f'<span style="font-size:{0.55*m:.2f}rem;color:white;font-weight:700;">{esc((mem.get("name","?") or "?")[0])}</span></div>'
            f'<div style="font-size:{0.52*m:.2f}rem;font-weight:700;color:{txt_color};">{esc(mem.get("name",""))}</div>'
            f'<div style="font-size:{0.48*m:.2f}rem;color:{acc};margin-top:1px;">{esc(mem.get("role",""))}</div></div>'
            for mem in members])
        inner = (
            f'{header_bar}'
            f'<div style="position:absolute;top:{content_top};left:0;right:0;bottom:0;'
            f'display:flex;align-items:center;justify-content:center;padding:0 6%;">'
            f'<div style="display:flex;gap:8px;width:100%;justify-content:center;">{m_html}</div></div>'
        )
    elif stype == 'two_column':
        bullets = slide.get('bullets', [])
        half = max(len(bullets) // 2, 1)
        lb, rb = bullets[:half], bullets[half:]
        col_bg = f'{acc}0D' if style != 'Dark' else f'{acc}22'
        def col_b(bl):
            return ''.join([
                f'<div style="display:flex;gap:3px;margin-bottom:3px;">'
                f'<span style="color:{acc};font-size:{0.55*m:.2f}rem;flex-shrink:0;">•</span>'
                f'<span style="font-size:{0.5*m:.2f}rem;color:{txt_color};opacity:0.85;">{esc(b)}</span></div>'
                for b in bl[:3]])
        inner = (
            f'{header_bar}'
            f'<div style="position:absolute;top:{content_top};left:0;right:0;bottom:0;padding:3% 5%;display:flex;gap:4%;">'
            f'<div style="flex:1;background:{col_bg};border-radius:5px;padding:5px;">{col_b(lb)}</div>'
            f'<div style="flex:1;background:{col_bg};border-radius:5px;padding:5px;">{col_b(rb)}</div></div>'
        )
    else:  # title_and_content
        bullets = slide.get('bullets', [])[:4]
        b_html = ''.join([
            f'<div style="display:flex;align-items:flex-start;gap:4px;margin-bottom:4px;">'
            f'{bullet_dot}'
            f'<div style="font-size:{0.53*m:.2f}rem;color:{txt_color};opacity:0.88;line-height:1.4;">{esc(b)}</div></div>'
            for b in bullets])
        inner = f'{header_bar}<div style="position:absolute;top:{content_top};left:0;right:0;bottom:0;padding:3% 6%;">{b_html}</div>'

    return (
        f'<div style="width:100%;padding-top:56.25%;position:relative;background:{actual_bg};'
        f'border-radius:10px;overflow:hidden;box-shadow:0 3px 15px rgba(0,0,0,0.15);'
        f'border:1px solid rgba(128,128,128,0.12);">'
        f'<div style="position:absolute;inset:0;">{inner}</div></div>'
    )


def _render_document_view_html(data, acc='#4F46E5', txt='#0F0F2A'):
    """Snapdeck3 스타일 문서 아웃라인 뷰 HTML 생성."""
    def esc(s):
        return str(s or '').replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')

    slides    = data.get('slides', [])
    pres_title = esc(data.get('presentation_title', '프레젠테이션'))
    pres_sub   = esc(data.get('presentation_subtitle', ''))

    html = (
        f'<div style="background:white;border-radius:16px;padding:36px 44px;'
        f'box-shadow:0 1px 3px rgba(0,0,0,0.06),0 4px 18px rgba(0,0,0,0.05);'
        f'font-family:inherit;line-height:1.6;">'
        f'<div style="margin-bottom:24px;">'
        f'<h1 style="font-size:1.45rem;font-weight:800;color:{txt};margin:0 0 5px;line-height:1.25;">{pres_title}</h1>'
        f'<p style="color:#aaa;font-size:0.82rem;margin:0;">{pres_sub}</p></div>'
        f'<hr style="border:none;border-top:2px solid #f0f0f0;margin:0 0 28px;">'
    )

    _TYPE_LABEL = {
        'title_and_content':'Content','section_header':'Section','big_stat':'Big Stat',
        'two_column':'Two Column','three_cards':'Three Cards','timeline':'Timeline','team_grid':'Team',
    }

    for i, slide in enumerate(slides):
        stype   = slide.get('slide_type','title_and_content')
        title   = esc(slide.get('title',''))
        bullets = [esc(b) for b in slide.get('bullets', [])]
        kw      = esc(slide.get('key_takeaway',''))
        summary = esc(slide.get('summary',''))
        tlabel  = _TYPE_LABEL.get(stype, stype.replace('_',' '))

        html += (
            f'<div id="sd-{i+1}" style="margin-bottom:28px;">'
            f'<div style="font-size:0.56rem;font-weight:700;letter-spacing:0.14em;'
            f'color:#ccc;margin-bottom:8px;text-transform:uppercase;">SLIDE {i+1} · {tlabel}</div>'
        )

        if stype == 'section_header':
            html += (
                f'<div style="background:{acc}0C;border-left:4px solid {acc};'
                f'border-radius:0 10px 10px 0;padding:16px 20px;">'
                f'<h2 style="font-size:1.3rem;font-weight:800;color:{acc};margin:0 0 5px;line-height:1.2;">{title}</h2>'
                f'<p style="font-size:0.85rem;color:#666;margin:0;">{summary}</p></div>'
            )
        elif stype == 'big_stat':
            sv = esc(slide.get('stat_value',''))
            sd = esc(slide.get('stat_description',''))
            html += (
                f'<h2 style="font-size:1.05rem;font-weight:600;color:{txt};margin:0 0 12px;">{title}</h2>'
                f'<div style="display:flex;align-items:baseline;gap:10px;margin-bottom:8px;">'
                f'<span style="font-size:2.6rem;font-weight:900;color:{acc};line-height:1;">{sv}</span>'
                f'<span style="font-size:0.88rem;color:#777;">{sd}</span></div>'
            )
            for b in bullets:
                html += f'<div style="display:flex;gap:8px;padding:3px 0;font-size:0.87rem;color:#555;"><span style="color:{acc};flex-shrink:0;">•</span>{b}</div>'
        elif stype == 'three_cards':
            cards = slide.get('cards', [])
            html += f'<h2 style="font-size:1.05rem;font-weight:600;color:{txt};margin:0 0 12px;">{title}</h2>'
            html += f'<div style="display:grid;grid-template-columns:repeat(3,1fr);gap:10px;">'
            for c in cards[:3]:
                ct = esc(c.get('card_title','')); cc = esc(c.get('card_content',''))
                html += (
                    f'<div style="background:{acc}07;border-radius:10px;padding:12px 14px;'
                    f'border:1px solid {acc}15;">'
                    f'<div style="font-size:0.8rem;font-weight:700;color:{acc};margin-bottom:5px;">{ct}</div>'
                    f'<div style="font-size:0.76rem;color:#555;line-height:1.5;">{cc}</div></div>'
                )
            html += '</div>'
        elif stype == 'timeline':
            tsteps = slide.get('timeline_steps', [])
            html += f'<h2 style="font-size:1.05rem;font-weight:600;color:{txt};margin:0 0 12px;">{title}</h2>'
            for j, s in enumerate(tsteps[:4]):
                st_t = esc(s.get('step_title','')); st_d = esc(s.get('step_desc',''))
                html += (
                    f'<div style="display:flex;gap:12px;margin-bottom:10px;align-items:flex-start;">'
                    f'<div style="width:22px;height:22px;border-radius:50%;background:{acc};flex-shrink:0;'
                    f'display:flex;align-items:center;justify-content:center;'
                    f'font-size:0.6rem;color:white;font-weight:800;margin-top:2px;">{j+1}</div>'
                    f'<div><div style="font-size:0.85rem;font-weight:600;color:{txt};">{st_t}</div>'
                    f'<div style="font-size:0.78rem;color:#777;">{st_d}</div></div></div>'
                )
        elif stype == 'two_column':
            half = max(len(bullets)//2, 1)
            html += f'<h2 style="font-size:1.05rem;font-weight:600;color:{txt};margin:0 0 12px;">{title}</h2>'
            html += '<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;">'
            for col_buls in [bullets[:half], bullets[half:]]:
                html += '<div style="background:#fafafa;border-radius:8px;padding:12px;">'
                for b in col_buls:
                    html += f'<div style="display:flex;gap:8px;padding:3px 0;font-size:0.85rem;color:#555;"><span style="color:{acc};">•</span>{b}</div>'
                html += '</div>'
            html += '</div>'
        elif stype == 'team_grid':
            members = slide.get('team_members', [])
            html += f'<h2 style="font-size:1.05rem;font-weight:600;color:{txt};margin:0 0 12px;">{title}</h2>'
            html += '<div style="display:flex;gap:14px;flex-wrap:wrap;">'
            for mem in members:
                n = esc(mem.get('name','')); r = esc(mem.get('role',''))
                init = n[0] if n else '?'
                html += (
                    f'<div style="text-align:center;min-width:72px;">'
                    f'<div style="width:38px;height:38px;border-radius:50%;background:{acc};margin:0 auto 6px;'
                    f'display:flex;align-items:center;justify-content:center;'
                    f'font-size:0.95rem;color:white;font-weight:700;">{init}</div>'
                    f'<div style="font-size:0.8rem;font-weight:600;color:{txt};">{n}</div>'
                    f'<div style="font-size:0.7rem;color:#aaa;">{r}</div></div>'
                )
            html += '</div>'
        else:  # title_and_content
            html += f'<h2 style="font-size:1.05rem;font-weight:600;color:{txt};margin:0 0 12px;line-height:1.3;">{title}</h2>'
            for b in bullets:
                html += (
                    f'<div style="display:flex;gap:10px;padding:4px 0;font-size:0.88rem;color:#444;line-height:1.55;">'
                    f'<span style="color:{acc};flex-shrink:0;font-size:0.6rem;margin-top:6px;">●</span>{b}</div>'
                )

        if kw:
            html += (
                f'<div style="margin-top:12px;padding:8px 14px;background:{acc}09;'
                f'border-radius:6px;font-size:0.78rem;color:#777;border-left:3px solid {acc};">'
                f'💡 {kw}</div>'
            )
        html += '</div>'
        if i < len(slides) - 1:
            html += '<hr style="border:none;border-top:1px solid #f4f4f4;margin:0 0 24px;">'

    html += '</div>'
    return html


def run_chat_edit(presentation_data, instruction, lang):
    """채팅 지시사항으로 특정 슬라이드(또는 전체)를 수정. (updated_data, response_msg) 반환."""
    slides = presentation_data.get('slides', [])

    # 슬라이드 번호 파싱
    nums = re.findall(r'(?:슬라이드|slide)\s*(\d+)', instruction, re.IGNORECASE)
    nums += re.findall(r'(\d+)\s*(?:번|번째)?\s*(?:슬라이드|slide|장)', instruction, re.IGNORECASE)
    nums = list(dict.fromkeys(nums))  # dedup

    if nums:
        target_idxs = [int(n)-1 for n in nums if 1 <= int(n) <= len(slides)]
    else:
        is_all = any(k in instruction for k in ['전체','모든','전부','다','all','every','모두'])
        if is_all:
            target_idxs = list(range(len(slides)))
        else:
            return None, "어떤 슬라이드를 수정할까요? 슬라이드 번호를 포함해 주세요.\n예) *슬라이드 3 제목을 더 강렬하게* / *전체 톤을 더 포멀하게*"

    modified = []
    for idx in target_idxs:
        if 0 <= idx < len(slides):
            new_sl = magic_edit_slide(slides[idx], instruction, lang)
            if new_sl:
                new_sl['_id'] = slides[idx].get('_id', str(idx))
                slides[idx] = new_sl
                modified.append(idx + 1)

    if not modified:
        return None, "수정에 실패했습니다. 다시 시도해 주세요."

    presentation_data['slides'] = slides
    return presentation_data, f"슬라이드 {', '.join(str(n) for n in modified)} 수정 완료 ✓"


def trigger_pptx_recompile():
    data = st.session_state.get('presentation_data')
    if data:
        buf = io.BytesIO()
        compile_presentation(
            data, buf,
            theme_name=st.session_state.get('theme'),
            deck_format=st.session_state.get('deck_format', 'Standard (16:9)'),
            template_style=st.session_state.get('template_style', 'Modern')
        )
        buf.seek(0)
        st.session_state['compiled_pptx'] = buf.read()


def render_progress_header():
    step = st.session_state.get('step', 1)
    steps = [(1,"프롬프트"), (2,"스타일"), (3,"AI 편집")]
    parts = []
    for i, (n, label) in enumerate(steps):
        cls = "done" if n < step else ("active" if n == step else "")
        icon = "✓" if n < step else str(n)
        parts.append(f'<div class="w-step {cls}"><span class="w-num">{icon}</span><span class="w-label">{label}</span></div>')
        if i < len(steps) - 1:
            lc = "done" if n < step else ""
            parts.append(f'<div class="w-line {lc}"></div>')
    st.markdown(f'<div class="wiz-bar">{"".join(parts)}</div>', unsafe_allow_html=True)


def render_navbar():
    dark = st.session_state.get('dark_mode', False)
    n1, n2, n3 = st.columns([2.5, 5, 1.5])
    with n1:
        st.markdown('<div class="nav-logo" style="padding:0.9rem 0;">⚡ snapdeck</div>', unsafe_allow_html=True)
    with n2:
        st.markdown('''
        <div class="nav-links" style="padding:0.85rem 0; justify-content:center;">
            <a href="#" class="nav-link">기능</a>
            <a href="#" class="nav-link">요금제</a>
            <a href="#" class="nav-link">블로그</a>
        </div>''', unsafe_allow_html=True)
    with n3:
        toggle_label = "🌙 다크" if not dark else "☀️ 라이트"
        if st.button(toggle_label, key="dark_toggle", use_container_width=True):
            st.session_state['dark_mode'] = not dark
            st.rerun()
    st.markdown('<hr style="margin:0; border:none; border-top:1px solid var(--border);">', unsafe_allow_html=True)


def _slide_form_key(slide, suffix):
    return f"{suffix}_{slide.get('_id', slide.get('slide_index', '0'))}"


# ════════════════════════════════════════════════════════════════
#  VIEW 1 · LANDING
# ════════════════════════════════════════════════════════════════
if st.session_state['view'] == 'landing':
    render_navbar()

    # ── Hero ──────────────────────────────────────────────────────
    st.markdown("""
    <div class="hero-outer">
        <div class="hero-glow"></div>
        <div class="hero-inner">
            <div class="hero-badge">✦ AI-Powered Presentation Builder</div>
            <div class="hero-title">
                발표 자료,<br><span class="g">이제 60초면 충분해요.</span>
            </div>
            <div class="hero-sub">
                주제를 입력하면 AI가 슬라이드 기획부터 디자인까지 자동으로 완성합니다.<br>
                이미지가 아닌 <strong style="color:var(--text);">완전히 편집 가능한 .pptx 파일</strong>로 바로 받아가세요.
            </div>
            <div class="hero-proof">
                <div class="hero-proof-item"><b>2,400+</b> 개 덱 생성</div>
                <div class="hero-proof-sep"></div>
                <div class="hero-proof-item"><b>4.9★</b> 사용자 평점</div>
                <div class="hero-proof-sep"></div>
                <div class="hero-proof-item">완전 <b>무료</b>로 시작</div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    h1, h2, h3 = st.columns([3, 2.2, 3])
    with h2:
        if st.button("무료로 시작하기  →", use_container_width=True, key="hero_cta"):
            st.session_state['view'] = 'signup'
            st.rerun()

    # ── Stats Strip ───────────────────────────────────────────────
    st.markdown("""
    <div class="stats-strip">
        <div class="stat-item">
            <div class="stat-num">50,000+</div>
            <div class="stat-label">생성된 슬라이드 수</div>
        </div>
        <div class="stat-item">
            <div class="stat-num">60초</div>
            <div class="stat-label">평균 덱 완성 시간</div>
        </div>
        <div class="stat-item">
            <div class="stat-num">7가지</div>
            <div class="stat-label">스마트 레이아웃 타입</div>
        </div>
        <div class="stat-item">
            <div class="stat-num">20개</div>
            <div class="stat-label">지원 언어</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── How It Works ──────────────────────────────────────────────
    st.markdown("""
    <div class="slabel">사용 방법</div>
    <div class="stitle">3단계로 끝납니다</div>
    <div class="ssub">복잡한 디자인 작업 없이, AI가 기획부터 파일 생성까지 전부 처리해드립니다.</div>
    <div class="steps-grid">
        <div class="step-card">
            <div class="step-num">01</div>
            <div class="step-title">주제와 목적을 입력하세요</div>
            <div class="step-desc">발표 주제, 언어, 목적(IR 덱 · 보고서 · 학습 발표 등), 슬라이드 수를 설정합니다. 참고 자료도 바로 붙여넣을 수 있어요.</div>
        </div>
        <div class="step-card">
            <div class="step-num">02</div>
            <div class="step-title">AI가 아웃라인을 기획합니다</div>
            <div class="step-desc">AI가 슬라이드 구조, 핵심 내용, 발표자 스크립트까지 자동 생성합니다. 내용은 직접 수정 · 추가 · 순서 변경도 자유롭습니다.</div>
        </div>
        <div class="step-card">
            <div class="step-num">03</div>
            <div class="step-title">편집 가능한 PPTX로 받아가세요</div>
            <div class="step-desc">이미지 캡처가 아닌 <strong style="color:var(--accent-l);">진짜 PowerPoint 파일</strong>로 내보냅니다. 글씨 수정, 색상 변경, 요소 이동이 완전히 자유롭습니다.</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Demo Mockup ───────────────────────────────────────────────
    st.markdown("""
    <div class="slabel">미리보기</div>
    <div class="stitle">이런 슬라이드가 만들어집니다</div>
    <div class="ssub">AI가 생성한 실제 출력 결과물입니다. 모든 텍스트와 요소는 PowerPoint에서 바로 수정 가능합니다.</div>
    <div class="mockup">
        <div class="mockup-bar">
            <div class="mdots">
                <div class="md md-r"></div>
                <div class="md md-y"></div>
                <div class="md md-g"></div>
            </div>
            <div class="mfile">snapdeck_output.pptx</div>
            <div class="msize">✦ 완전 편집 가능</div>
        </div>
        <div class="mockup-body">
            <div class="m-slide-hdr">
                <div class="m-title">AI 진단 보조 솔루션 — 투자 유치 덱</div>
                <div class="m-sub">의료 효율성 극대화와 환자 안전을 위한 AI 혁신 · 시리즈 A</div>
            </div>
            <div class="m-cards">
                <div class="m-card">
                    <div class="m-ctitle">⚡ EHR 직접 연동 API</div>
                    <div class="m-cdesc">기존 병원 시스템에 별도 인프라 없이 48시간 내 배치 완료.</div>
                </div>
                <div class="m-card">
                    <div class="m-ctitle">📊 진단 정확도 94.7%</div>
                    <div class="m-cdesc">3개 상급병원 임상 검증 완료. FDA 510(k) 준비 중.</div>
                </div>
                <div class="m-card">
                    <div class="m-ctitle">💰 ARR $2.1M</div>
                    <div class="m-cdesc">전년 대비 340% 성장. 파이프라인 $8M 추가 확보.</div>
                </div>
            </div>
            <div class="m-tw">💡 Key Takeaway: 환자 안전을 보장하는 검증된 AI 어시스턴트 — 병원 수익성도 동시에 개선</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Features ──────────────────────────────────────────────────
    st.markdown("""
    <div class="slabel">핵심 기능</div>
    <div class="stitle">왜 Snapdeck인가요?</div>
    <div class="ssub">단순한 슬라이드 생성을 넘어, 기획부터 파일 완성까지 발표의 전 과정을 AI가 처리합니다.</div>
    <div class="feat-grid">
        <div class="feat">
            <span class="feat-icon">🤖</span>
            <div class="feat-title">AI 아웃라인 자동 기획</div>
            <div class="feat-desc">주제만 입력하면 슬라이드 구조, 흐름, 핵심 메시지를 AI가 즉시 기획합니다. 빈 화면에서 시작할 필요가 없어요.</div>
        </div>
        <div class="feat">
            <span class="feat-icon">✨</span>
            <div class="feat-title">Magic Edit</div>
            <div class="feat-desc">자연어로 명령하면 내용을 즉시 개선합니다. "임원 발표용으로 다듬어줘", "영어로 바꿔줘" 한 마디면 완성.</div>
        </div>
        <div class="feat">
            <span class="feat-icon">📐</span>
            <div class="feat-title">7가지 스마트 레이아웃</div>
            <div class="feat-desc">카드 · 타임라인 · 통계 · 팀 그리드 등 내용에 맞는 레이아웃이 자동으로 선택됩니다. 직접 변경도 가능합니다.</div>
        </div>
        <div class="feat">
            <span class="feat-icon">📥</span>
            <div class="feat-title">진짜 .pptx 파일로 내보내기</div>
            <div class="feat-desc">이미지 캡처가 아닌 완전한 PowerPoint 파일로 다운로드합니다. 글씨, 색상, 도형 — 모두 편집 가능합니다.</div>
        </div>
        <div class="feat">
            <span class="feat-icon">🗒️</span>
            <div class="feat-title">발표자 스크립트 자동 생성</div>
            <div class="feat-desc">각 슬라이드마다 발표자 노트를 AI가 자동으로 작성합니다. 발표 준비 시간이 절반으로 줄어듭니다.</div>
        </div>
        <div class="feat">
            <span class="feat-icon">🌐</span>
            <div class="feat-title">20개 언어 지원</div>
            <div class="feat-desc">한국어 · 영어 · 일본어 · 중국어 등 20개 언어로 발표 자료를 즉시 생성합니다. 글로벌 발표도 문제없어요.</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Pricing ───────────────────────────────────────────────────
    st.markdown("""
    <div class="slabel">요금제</div>
    <div class="stitle">심플한 가격</div>
    <div class="ssub">숨겨진 비용 없이, 필요한 만큼만 사용하세요. 지금 당장 무료로 시작하세요.</div>
    <div class="price-wrap">
        <div class="plan">
            <div class="plan-name">Free</div>
            <div class="plan-price">$0</div>
            <ul class="plan-list">
                <li>AI 슬라이드 기본 생성</li>
                <li>기본 레이아웃 4종</li>
                <li>PPTX 내보내기 (워터마크)</li>
                <li>월 5회 생성 제한</li>
            </ul>
        </div>
        <div class="plan pro">
            <div class="plan-badge">✦ 추천</div>
            <div class="plan-name">Pro</div>
            <div class="plan-price">$15 <span class="plan-unit">/ 월</span></div>
            <ul class="plan-list">
                <li>무제한 AI 슬라이드 생성</li>
                <li>모든 레이아웃 · 테마 4종</li>
                <li>워터마크 없는 PPTX</li>
                <li>Magic Edit 무제한</li>
                <li>발표자 스크립트 자동 생성</li>
                <li>브랜드 키트 (로고 · 색상)</li>
            </ul>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── CTA Banner ────────────────────────────────────────────────
    st.markdown("""
    <div class="cta-banner">
        <div class="cta-banner-inner">
            <div class="cta-banner-title">지금 바로 첫 번째 덱을 만들어보세요.</div>
            <div class="cta-banner-sub">가입 후 60초 안에 첫 번째 슬라이드를 받아볼 수 있습니다.<br>신용카드 불필요 · 무료로 시작</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    cp1, cp2, cp3, cp4 = st.columns([2, 2, 2, 2])
    with cp2:
        if st.button("무료로 시작하기", key="cta_free", use_container_width=True):
            st.session_state['view'] = 'signup'; st.rerun()
    with cp3:
        if st.button("Pro 업그레이드", key="cta_pro", use_container_width=True):
            st.session_state['view'] = 'signup'; st.rerun()


# ════════════════════════════════════════════════════════════════
#  VIEW 2 · SIGNUP
# ════════════════════════════════════════════════════════════════
elif st.session_state['view'] == 'signup':
    render_navbar()
    st.markdown("""
    <div class="auth-box">
        <div class="auth-title">계정 만들기</div>
        <div class="auth-sub">무료로 시작하고, 언제든 Pro로 업그레이드하세요.</div>
    </div>
    """, unsafe_allow_html=True)

    sa1, sa2, sa3 = st.columns([3, 4, 3])
    with sa2:
        if st.button("🔴  Google 계정으로 계속하기", key="btn_g", use_container_width=True):
            st.session_state['view'] = 'onboarding'; st.rerun()
        if st.button("💛  카카오 계정으로 계속하기", key="btn_k", use_container_width=True):
            st.session_state['view'] = 'onboarding'; st.rerun()
        if st.button("🍎  Apple 계정으로 계속하기", key="btn_a", use_container_width=True):
            st.session_state['view'] = 'onboarding'; st.rerun()

        st.markdown('<div class="divline">또는 이메일로 가입</div>', unsafe_allow_html=True)
        email = st.text_input("이메일 주소", placeholder="name@company.com")
        pw = st.text_input("비밀번호", type="password")
        st.write("")
        if st.button("회원가입", key="btn_email", use_container_width=True):
            if email and pw:
                st.session_state['view'] = 'onboarding'; st.rerun()
            else:
                st.warning("이메일과 비밀번호를 입력해 주세요.")
        st.markdown("<div style='text-align:center;margin-top:1rem;'><a href='#' style='color:#A78BFA;font-size:0.82rem;text-decoration:none;'>이미 계정이 있으신가요? 로그인</a></div>", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════
#  VIEW 3 · ONBOARDING
# ════════════════════════════════════════════════════════════════
elif st.session_state['view'] == 'onboarding':
    render_navbar()
    st.markdown("## Snapdeck 온보딩")
    st.markdown("작업 방식을 알려주시면 최적의 템플릿과 스타일을 제안해 드립니다.")
    st.divider()

    ob1, ob2, ob3 = st.columns([1, 4, 1])
    with ob2:
        st.markdown("#### 주로 어떤 언어로 작업하시나요?")
        st.caption("선택한 언어를 기준으로 슬라이드를 기획합니다.")
        langs_all = ["한국어","English","日本語","中文 (简体)","中文 (繁體)",
                     "Español","Français","Deutsch","Português","Italiano",
                     "Русский","العربية","हिंदी","Nederlands","Polski",
                     "Türkçe","Tiếng Việt","ภาษาไทย","Bahasa","Svenska"]
        lc = st.columns(4)
        sel_lang = "한국어"
        for idx, lg in enumerate(langs_all):
            with lc[idx % 4]:
                if st.checkbox(lg, value=(lg == "한국어"), key=f"ob_l_{lg}"):
                    sel_lang = lg

        st.write("")
        st.markdown("#### 어떤 유형의 발표를 주로 만드시나요?")
        cats = [
            "Startup Pitch (투자 유치 및 IR 덱)",
            "Business Report (전략 · 분석 · 인사이트)",
            "Academic (연구 및 학습 발표)",
            "Sales Deck (제안서 및 세일즈 덱)",
            "Other (그 외 목적)",
        ]
        sel_cats = []
        for cat in cats:
            if st.checkbox(cat, value=(cat == cats[0]), key=f"ob_c_{cat}"):
                sel_cats.append(cat)

        st.session_state['lang'] = sel_lang
        if sel_cats:
            st.session_state['category'] = sel_cats[0]

        st.write("")
        if st.button("서비스 시작하기  →", key="ob_submit", use_container_width=True):
            st.session_state['view'] = 'generator'; st.rerun()


# ════════════════════════════════════════════════════════════════
#  VIEW 4 · GENERATOR
# ════════════════════════════════════════════════════════════════
elif st.session_state['view'] == 'generator':

    # Sidebar
    with st.sidebar:
        st.markdown("### ⚡ snapdeck")
        st.divider()
        if not api_key:
            st.warning("API Key가 설정되지 않았습니다.")
            entered = st.text_input("Gemini API Key", type="password", key="key_input")
            if entered:
                st.session_state['gemini_api_key'] = entered
                with open(".env", "w", encoding="utf-8") as f:
                    f.write(f"GEMINI_API_KEY={entered}\n")
                st.success("저장 완료!")
                st.rerun()
        else:
            st.success("API Key 연결됨")
            if st.button("API Key 변경"):
                st.session_state['gemini_api_key'] = ''
                if os.path.exists(".env"):
                    os.remove(".env")
                st.rerun()
        st.divider()
        if st.button("← 메인으로 돌아가기"):
            st.session_state['view'] = 'landing'
            st.session_state['step'] = 1
            st.rerun()

    render_progress_header()

    # ─── STEP 1: 프롬프트 입력 ───────────────────────────────────────
    if st.session_state['step'] == 1:
        st.markdown("### 📝 발표 주제 및 정보 입력")
        st.markdown("만들고 싶은 발표 자료를 자세히 설명할수록 더 완성도 높은 슬라이드가 생성됩니다.")
        st.divider()

        topic = st.text_area(
            "발표 주제 및 핵심 기획 아이디어",
            value=st.session_state['topic'],
            placeholder="예: 헬스케어 비효율 해소를 위한 AI 진단 솔루션 투자 유치 자료. 문제점, 솔루션 특징, ARR 실적, 팀 구성, 자금 조달 계획 포함.",
            height=110,
        )

        c_l, c_r = st.columns(2)
        with c_l:
            slide_count = st.slider(
                "슬라이드 수",
                min_value=3, max_value=20,
                value=st.session_state['slide_count'],
                help="총 생성할 슬라이드 장 수 (타이틀 슬라이드 포함)"
            )
        with c_r:
            langs_list = ["한국어","English","日本語","中文 (简体)","中文 (繁體)",
                          "Español","Français","Deutsch","Português","Italiano",
                          "Русский","العربية","हिंदी","Nederlands","Polski",
                          "Türkçe","Tiếng Việt","ภาษาไทย","Bahasa","Svenska"]
            lang = st.selectbox("언어", langs_list,
                                index=langs_list.index(st.session_state['lang']))

        cats_list = [
            "Startup Pitch (투자 유치 및 IR 덱)",
            "Business Report (전략 · 분석 · 인사이트)",
            "Academic (연구 및 학습 발표)",
            "Sales Deck (제안서 및 세일즈 덱)",
            "Other (그 외 목적)",
        ]
        category = st.selectbox("발표 목적", cats_list,
                                index=cats_list.index(st.session_state['category']))

        st.divider()
        st.markdown("#### 참고 자료 추가 (선택)")
        st.caption("기획서, 보고서, 원본 텍스트, PPTX 파일, 또는 웹 URL을 제공하면 AI가 내용을 분석해 슬라이드에 반영합니다.")

        src_tab1, src_tab2, src_tab3 = st.tabs(["📄 파일 업로드", "🌐 URL 가져오기", "✏️ 직접 입력"])

        with src_tab1:
            up = st.file_uploader("파일 업로드 (.txt, .md, .pdf, .pptx)", type=["txt","md","pdf","pptx"])
            if up is not None:
                if st.session_state.get('_upload_name') != up.name:
                    try:
                        if up.name.lower().endswith(".pdf"):
                            reader = PdfReader(up)
                            text = "\n".join(page.extract_text() or "" for page in reader.pages)
                            st.session_state['source_text'] = text.strip()
                        elif up.name.lower().endswith(".pptx"):
                            text = extract_pptx_text(up)
                            st.session_state['source_text'] = text
                        else:
                            st.session_state['source_text'] = up.read().decode("utf-8")
                        st.session_state['_upload_name'] = up.name
                    except Exception as e:
                        st.error(f"파일 오류: {e}")
                char_count = len(st.session_state.get('source_text', ''))
                st.success(f"✓ {up.name} 로드 완료 ({char_count:,}자)")
                if st.session_state.get('source_text'):
                    with st.expander("추출된 내용 미리보기"):
                        st.text(st.session_state['source_text'][:800] + ("..." if len(st.session_state.get('source_text','')) > 800 else ""))
            else:
                st.session_state['_upload_name'] = None

        with src_tab2:
            url_input = st.text_input(
                "웹 페이지 URL",
                placeholder="https://example.com/article",
                key='_url_input',
            )
            if st.button("URL 내용 가져오기", key='_fetch_url'):
                if url_input.strip():
                    with st.spinner("웹 페이지 분석 중..."):
                        text, err = fetch_url_text(url_input)
                    if err:
                        st.error(f"가져오기 실패: {err}")
                    else:
                        st.session_state['source_text'] = text
                        st.session_state['_url_loaded'] = url_input
                        st.success(f"✓ {len(text):,}자 추출 완료")
                        st.rerun()
            if st.session_state.get('_url_loaded'):
                st.caption(f"로드됨: {st.session_state['_url_loaded']}")
                with st.expander("추출된 내용 미리보기"):
                    txt = st.session_state.get('source_text', '')
                    st.text(txt[:800] + ("..." if len(txt) > 800 else ""))

        with src_tab3:
            st.text_area(
                "텍스트 직접 붙여넣기",
                key='source_text',
                placeholder="노션, 메일, 문서 등에서 복사한 내용을 여기에 붙여넣으세요...",
                height=150,
            )

        st.write("")
        if st.button("다음 단계 — 스타일 선택  →", use_container_width=False):
            if not topic.strip():
                st.warning("발표 주제를 먼저 입력해 주세요.")
            else:
                st.session_state['topic'] = topic
                st.session_state['slide_count'] = slide_count
                st.session_state['lang'] = lang
                st.session_state['category'] = category
                st.session_state['step'] = 2
                st.rerun()

    # ─── STEP 2: 스타일 & 테마 ───────────────────────────────────────
    elif st.session_state['step'] == 2:

        THEME_META = [
            {"key":"Nordic Tech (Ice White & Indigo)",        "short":"노딕 테크",    "tag":"IR · 스타트업",  "desc":"깔끔하고 신뢰감 있는 테크 스타일. IR·투자 유치 피치에 최적."},
            {"key":"Obsidian Dark (Slate Dark & Cyan)",       "short":"옵시디안 다크","tag":"테크 · 개발",    "desc":"다크 배경에 시안 강조. 기술 제품·개발 발표에 어울립니다."},
            {"key":"Emerald Business (Mint & Forest)",        "short":"에메랄드",     "tag":"비즈니스 · 전략","desc":"민트·포레스트 그린 조합. 안정적인 기업 발표에 적합합니다."},
            {"key":"Sunset Warm (Cream & Terracotta)",        "short":"선셋 웜",      "tag":"브랜드 · 마케팅","desc":"따뜻한 크림·테라코타. 브랜드와 마케팅 덱에 어울립니다."},
            {"key":"Midnight Purple (Deep Violet & Lavender)","short":"미드나잇",     "tag":"크리에이티브",   "desc":"딥 퍼플·라벤더. 크리에이티브 에이전시·디자인 발표."},
            {"key":"Ocean Blue (Navy & Sky)",                 "short":"오션 블루",    "tag":"금융 · 리포트",  "desc":"네이비·스카이블루 조합. 금융·리포트 발표에 최적입니다."},
            {"key":"Rose Gold (Ivory & Blush)",               "short":"로즈 골드",    "tag":"럭셔리 · 패션",  "desc":"아이보리·블러쉬 골드. 럭셔리·패션 브랜드 발표."},
            {"key":"Monochrome (Pure White & Black)",         "short":"모노크롬",     "tag":"미니멀",         "desc":"퓨어 화이트·블랙. 군더더기 없는 미니멀 발표."},
            {"key":"Cherry Blossom (Soft Pink & Rose)",       "short":"체리 블로섬",  "tag":"소셜 · 라이프",  "desc":"소프트 핑크·로즈. 소셜·라이프스타일 콘텐츠에 어울립니다."},
            {"key":"Custom",                                   "short":"커스텀",       "tag":"직접 설정",      "desc":"강조색·배경색·텍스트색을 직접 지정해 나만의 테마를 만듭니다."},
        ]
        CAT_REC = {
            "Startup Pitch":   "Nordic Tech (Ice White & Indigo)",
            "Business Report": "Ocean Blue (Navy & Sky)",
            "Academic":        "Emerald Business (Mint & Forest)",
            "Sales Deck":      "Sunset Warm (Cream & Terracotta)",
        }
        cat_key   = st.session_state.get('category','')
        rec_theme = next((v for k,v in CAT_REC.items() if k in cat_key), None)
        FMT_META  = [
            {"key":"Standard (16:9)",                  "w":16,  "h":9,   "ratio":"16:9",  "label":"스탠다드",    "desc":"일반 발표"},
            {"key":"Keynote Widescreen (16:10)",        "w":16,  "h":10,  "ratio":"16:10", "label":"와이드스크린","desc":"Mac·키노트"},
            {"key":"Classic Slide (4:3)",               "w":4,   "h":3,   "ratio":"4:3",   "label":"클래식",      "desc":"빔프로젝터"},
            {"key":"Social Carousel (4:5)",             "w":4,   "h":5,   "ratio":"4:5",   "label":"소셜 카드",   "desc":"인스타·SNS"},
            {"key":"Report / Document (A4-Horizontal)", "w":297, "h":210, "ratio":"A4-H",  "label":"리포트",      "desc":"문서·PDF"},
        ]
        sel_theme = st.session_state.get('theme','Nordic Tech (Ice White & Indigo)')
        sel_fmt   = st.session_state.get('deck_format','Standard (16:9)')
        ta, tb, tt = _get_colors(sel_theme)
        sel_tm = next((m for m in THEME_META if m['key']==sel_theme), THEME_META[0])
        sel_fm = next((m for m in FMT_META   if m['key']==sel_fmt),   FMT_META[0])

        # ── 헤더 ─────────────────────────────────────────────────────
        st.markdown("### 슬라이드 스타일 선택")
        st.caption("테마와 덱 비율을 선택하고, 색상을 원하는 대로 조정하세요.")
        st.divider()

        # ── 덱 비율 (상단 pill 행) ─────────────────────────────────
        st.markdown('<div style="font-size:0.78rem;font-weight:600;color:var(--text2);letter-spacing:0.04em;margin-bottom:8px;">덱 비율</div>', unsafe_allow_html=True)
        fc = st.columns(5)
        for ci, fm in enumerate(FMT_META):
            is_f = sel_fmt == fm['key']
            with fc[ci]:
                border = '#7C3AED' if is_f else 'rgba(99,102,241,0.15)'
                bg_f   = 'linear-gradient(135deg,#7C3AED,#4F46E5)' if is_f else 'var(--surface)'
                tc_f   = '#fff' if is_f else 'var(--text)'
                tc2_f  = 'rgba(255,255,255,0.75)' if is_f else 'var(--text2)'
                st.markdown(f"""
<div style="border:2px solid {border};border-radius:10px;padding:10px 6px 8px;
            background:{bg_f};text-align:center;margin-bottom:4px;
            box-shadow:{'0 4px 12px rgba(124,58,237,0.25)' if is_f else 'none'};">
  <div style="font-size:1rem;font-weight:800;color:{tc_f};line-height:1;">{fm['ratio']}</div>
  <div style="font-size:0.7rem;font-weight:600;color:{tc_f};margin-top:3px;">{fm['label']}</div>
  <div style="font-size:0.62rem;color:{tc2_f};margin-top:2px;">{fm['desc']}</div>
</div>""", unsafe_allow_html=True)
                lbl = "✓ 선택됨" if is_f else "선택"
                if st.button(lbl, key=f"fmt_{ci}", use_container_width=True):
                    st.session_state['deck_format'] = fm['key']
                    st.rerun()

        st.markdown('<div style="margin-top:10px;"></div>', unsafe_allow_html=True)
        st.divider()

        # ── 템플릿 스타일 ──────────────────────────────────────────
        STYLE_META = [
            {"key": "Modern",  "icon": "◼", "label": "모던",   "desc": "헤더 바 + 정돈된 불릿. 가장 범용적인 스타일."},
            {"key": "Bold",    "icon": "▮", "label": "볼드",   "desc": "강렬한 풀-와이드 헤더, 임팩트 있는 제목 강조."},
            {"key": "Minimal", "icon": "─", "label": "미니멀", "desc": "여백 중심, 깔끔한 액센트 타이틀. 전달력 최우선."},
            {"key": "Dark",    "icon": "◈", "label": "다크",   "desc": "다크 배경 + 강렬한 액센트. 세련된 테크 발표."},
        ]
        sel_style  = st.session_state.get('template_style', 'Modern')
        _style_dk  = _get_dark_bg_hex(sel_theme)
        st.markdown('<div style="font-size:0.78rem;font-weight:600;color:var(--text2);letter-spacing:0.04em;margin-bottom:8px;">레이아웃 스타일</div>', unsafe_allow_html=True)
        sc = st.columns(4)
        _demo_slide = {"slide_type": "title_and_content", "title": "슬라이드 제목",
                       "bullets": ["첫 번째 내용 포인트", "두 번째 내용 포인트", "세 번째 포인트"]}
        for ci, sm in enumerate(STYLE_META):
            is_sel = sel_style == sm["key"]
            with sc[ci]:
                border  = '#7C3AED' if is_sel else 'rgba(99,102,241,0.15)'
                bsh     = '0 4px 12px rgba(124,58,237,0.25)' if is_sel else 'none'
                hdr_clr = '#7C3AED' if is_sel else 'var(--text2)'
                st.markdown(
                    f'<div style="border:2px solid {border};border-radius:10px;padding:8px;'
                    f'margin-bottom:4px;box-shadow:{bsh};">'
                    f'<div style="font-size:1.1rem;text-align:center;margin-bottom:4px;color:{hdr_clr};">{sm["icon"]}</div>'
                    f'<div style="font-size:0.72rem;font-weight:700;text-align:center;color:var(--text);">{sm["label"]}</div>'
                    f'<div style="font-size:0.58rem;color:var(--text2);text-align:center;margin-top:2px;line-height:1.3;">{sm["desc"]}</div>'
                    f'<div style="margin-top:6px;">'
                    f'{_render_slide_preview_html(_demo_slide, ta, tb, tt, sm["key"], _style_dk)}'
                    f'</div></div>',
                    unsafe_allow_html=True
                )
                lbl = "✓ 선택됨" if is_sel else "선택"
                if st.button(lbl, key=f"sty_{sm['key']}", use_container_width=True):
                    st.session_state['template_style'] = sm['key']
                    st.rerun()

        st.markdown('<div style="margin-top:10px;"></div>', unsafe_allow_html=True)
        st.divider()

        # ── 메인 2-컬럼: 테마 리스트 | 프리뷰 ──────────────────────
        st.markdown('<div style="font-size:0.78rem;font-weight:600;color:var(--text2);letter-spacing:0.04em;margin-bottom:10px;">테마 선택</div>', unsafe_allow_html=True)

        list_col, prev_col = st.columns([4, 6], gap="large")

        with list_col:
            for tm in THEME_META:
                is_s = sel_theme == tm['key']
                is_r = tm['key'] == rec_theme
                a2, b2, t2 = _get_colors(tm['key'])
                border = '#7C3AED' if is_s else ('#F59E0B' if is_r else 'var(--border-b)')
                bg_row = 'var(--accent-bg)' if is_s else 'var(--surface)'
                ck_html = f'<span style="color:#7C3AED;font-size:1rem;font-weight:700;">✓</span>' if is_s else ''
                rec_html = '<span style="font-size:0.58rem;background:#FEF3C7;color:#92400E;padding:1px 5px;border-radius:5px;margin-left:4px;">추천</span>' if is_r else ''
                st.markdown(f"""
<div style="display:flex;align-items:center;gap:10px;padding:9px 10px;
            border-radius:12px;border:2px solid {border};background:{bg_row};margin-bottom:6px;">
  <div style="width:80px;flex-shrink:0;">
    {_slide_html(a2,b2,t2,16,9,5,False)}
  </div>
  <div style="flex:1;min-width:0;">
    <div style="font-size:0.83rem;font-weight:700;color:var(--text);
                display:flex;align-items:center;gap:4px;flex-wrap:wrap;">
      {tm['short']} {rec_html} {ck_html}
    </div>
    <div style="display:flex;gap:4px;margin-top:5px;align-items:center;">
      <div style="width:13px;height:13px;border-radius:3px;background:{b2};border:1px solid rgba(0,0,0,0.1);"></div>
      <div style="width:13px;height:13px;border-radius:3px;background:{a2};"></div>
      <div style="width:13px;height:13px;border-radius:3px;background:{t2};"></div>
      <span style="font-size:0.62rem;color:var(--text2);margin-left:3px;">{tm['tag']}</span>
    </div>
  </div>
</div>""", unsafe_allow_html=True)
                if not is_s:
                    if st.button("선택", key=f"th3_{tm['key'][:10]}", use_container_width=True):
                        st.session_state['theme'] = tm['key']
                        st.rerun()
                else:
                    st.markdown('<div style="height:32px;"></div>', unsafe_allow_html=True)

        with prev_col:
            # 큰 프리뷰 (덱 겹침 효과)
            st.markdown(f"""
<div style="position:relative;margin:8px 14px 0 0;">
  <div style="position:absolute;top:-7px;right:-7px;left:7px;bottom:7px;
              background:{tb};border-radius:12px;opacity:0.28;
              box-shadow:0 2px 8px rgba(0,0,0,0.1);"></div>
  <div style="position:absolute;top:-3px;right:-3px;left:3px;bottom:3px;
              background:{tb};border-radius:12px;opacity:0.6;
              box-shadow:0 2px 8px rgba(0,0,0,0.1);"></div>
  <div style="position:relative;">
    {_slide_html(ta, tb, tt, sel_fm['w'], sel_fm['h'], 10, True)}
  </div>
</div>""", unsafe_allow_html=True)

            # 테마 정보
            st.markdown(f"""
<div style="margin-top:24px;padding:0 2px;">
  <div style="display:flex;align-items:center;gap:8px;">
    <span style="font-size:1.08rem;font-weight:800;color:var(--text);">{sel_tm['short']}</span>
    <span style="font-size:0.68rem;background:var(--accent-bg);color:var(--accent);
                 padding:2px 8px;border-radius:8px;">{sel_tm['tag']}</span>
  </div>
  <p style="font-size:0.82rem;color:var(--text2);margin:6px 0 12px;">{sel_tm['desc']}</p>
  <div style="display:flex;gap:8px;align-items:center;">
    <div style="width:30px;height:30px;border-radius:8px;background:{tb};
                border:1.5px solid rgba(0,0,0,0.1);box-shadow:0 2px 6px rgba(0,0,0,0.08);"></div>
    <div style="width:30px;height:30px;border-radius:8px;background:{ta};
                box-shadow:0 2px 6px rgba(0,0,0,0.15);"></div>
    <div style="width:30px;height:30px;border-radius:8px;background:{tt};
                box-shadow:0 2px 6px rgba(0,0,0,0.08);"></div>
    <span style="font-size:0.7rem;color:var(--text2);margin-left:4px;">배경 · 강조 · 텍스트</span>
  </div>
</div>""", unsafe_allow_html=True)

            # 색상 커스터마이즈 — 선택된 테마 색상을 picker 기본값으로 동기화
            st.markdown('<div style="margin-top:18px;"></div>', unsafe_allow_html=True)
            _ck_init = '_color_init_from'
            _prev_ck = st.session_state.get(_ck_init, '')
            if sel_theme != 'Custom' and sel_theme != _prev_ck:
                _base = THEMES.get(sel_theme, list(THEMES.values())[0])
                st.session_state['custom_accent'] = "#{:02X}{:02X}{:02X}".format(*_base["accent"])
                st.session_state['custom_bg']     = "#{:02X}{:02X}{:02X}".format(*_base["bg_light"])
                st.session_state['custom_text']   = "#{:02X}{:02X}{:02X}".format(*_base["primary_text"])
                st.session_state[_ck_init] = sel_theme
            elif sel_theme == 'Custom' and _prev_ck != 'Custom':
                st.session_state[_ck_init] = 'Custom'

            with st.expander("색상 커스터마이즈", expanded=(sel_theme == "Custom")):
                pc1, pc2, pc3 = st.columns(3)
                with pc1: ca = st.color_picker("강조색", key='custom_accent')
                with pc2: cb = st.color_picker("배경색", key='custom_bg')
                with pc3: ct = st.color_picker("텍스트색", key='custom_text')
                # Always sync THEMES["Custom"] so the preview & PPTX stay current
                THEMES["Custom"]["accent"]         = RGBColor(*_hex_to_rgb(ca))
                THEMES["Custom"]["accent_light"]   = RGBColor(*_hex_to_rgb(ca))
                THEMES["Custom"]["accent_bg"]      = RGBColor(*_hex_to_rgb(cb))
                THEMES["Custom"]["bg_light"]       = RGBColor(*_hex_to_rgb(cb))
                THEMES["Custom"]["bg_dark"]        = RGBColor(*_hex_to_rgb(cb))
                THEMES["Custom"]["primary_text"]   = RGBColor(*_hex_to_rgb(ct))
                THEMES["Custom"]["secondary_text"] = RGBColor(*_hex_to_rgb(ct))
                st.markdown(_slide_html(ca, cb, ct, 16, 9, 8, True), unsafe_allow_html=True)
                if sel_theme != 'Custom':
                    if st.button("이 색상으로 커스텀 테마 적용", type="primary", use_container_width=True):
                        st.session_state['theme'] = 'Custom'
                        st.session_state[_ck_init] = 'Custom'
                        st.rerun()
                else:
                    st.success("✓ 커스텀 테마 적용 중 — 색상을 바꾸면 즉시 반영됩니다")

        st.divider()
        n1, n2, _ = st.columns([1.2, 1.2, 8])
        with n1:
            if st.button("← 이전"):
                st.session_state['step'] = 1; st.rerun()
        with n2:
            if st.button("다음 →", type="primary"):
                st.session_state['step'] = 3; st.rerun()

    # ─── STEP 3: AI 생성 + 문서 편집 (통합) ────────────────────────
    elif st.session_state['step'] == 3:
        _s3_theme = st.session_state.get('theme', 'Nordic Tech (Ice White & Indigo)')
        _s3_sty   = st.session_state.get('template_style', 'Modern')
        _s3_acc, _s3_bg, _s3_txt = _get_colors(_s3_theme)
        _s3_dark  = _get_dark_bg_hex(_s3_theme)

        data   = st.session_state.get('presentation_data') or {}
        slides = data.get('slides', [])
        lang   = st.session_state.get('lang', 'Korean')
        chat_history = st.session_state.get('chat_history', [])
        agent_stage  = st.session_state.get('agent_stage', 0)

        # ── 텍스트 인풋을 문서처럼 보이게 하는 CSS ─────────────────
        st.markdown(
            '<style>'
            '.doc-slide-title > div > div > input {'
            '  font-size:1.15rem;font-weight:700;border:none;border-bottom:1.5px solid #ebebf5;'
            '  border-radius:0;background:transparent;padding:4px 2px;color:#0F0F2A;}'
            '.doc-slide-title > div > div > input:focus {'
            '  border-bottom-color:#4F46E5;box-shadow:none;background:transparent;}'
            '.doc-bullet > div > div > input {'
            '  font-size:0.9rem;border:none;border-bottom:1px solid transparent;'
            '  border-radius:0;background:transparent;padding:3px 2px;color:#444;}'
            '.doc-bullet > div > div > input:hover {border-bottom-color:#e8e8f4;}'
            '.doc-bullet > div > div > input:focus {border-bottom-color:#4F46E5;box-shadow:none;background:transparent;}'
            '.doc-kw > div > div > input {'
            '  font-size:0.82rem;border:1px solid #ebebf5;border-radius:6px;'
            '  background:#fafafa;padding:4px 10px;color:#666;}'
            '</style>',
            unsafe_allow_html=True
        )

        # ── 상단 바 ──────────────────────────────────────────────────
        bar_a, bar_b, bar_c = st.columns([5, 2, 2])
        with bar_a:
            topic_lbl = st.session_state.get('topic', '')[:50]
            style_lbl = f"{_s3_theme.split('(')[0].strip()} · {_s3_sty}"
            slide_lbl = f" · {len(slides)}장" if slides else ''
            st.markdown(
                f'<div style="padding:8px 0 2px;">'
                f'<b style="font-size:1.05rem;">{topic_lbl}</b>'
                f'&ensp;<span style="font-size:0.78rem;color:#b0b0c0;">{style_lbl}{slide_lbl}</span>'
                f'</div>',
                unsafe_allow_html=True
            )
        with bar_b:
            if st.button('← 스타일', use_container_width=True):
                st.session_state['step'] = 2
                st.session_state['agent_stage'] = 0
                st.rerun()
        with bar_c:
            if slides and st.session_state.get('compiled_pptx'):
                _safe = re.sub(r'[\\/*?:"<>| ]', '_', st.session_state.get('topic', 'deck'))
                st.download_button(
                    '📥 PPTX 다운로드',
                    data=st.session_state['compiled_pptx'],
                    file_name=f'{_safe}.pptx',
                    mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    use_container_width=True,
                )

        st.markdown('<div style="height:4px;"></div>', unsafe_allow_html=True)

        if not slides:
            # ── 생성 전: Generate 화면 ─────────────────────────────
            st.markdown('<div style="height:60px;"></div>', unsafe_allow_html=True)
            _, gen_col, _ = st.columns([2, 4, 2])
            with gen_col:
                use_polish = st.checkbox('품질 에이전트 사용 (느리지만 문장 퀄리티 ↑)', value=False, key='use_polish_s3')
                st.markdown('<div style="height:12px;"></div>', unsafe_allow_html=True)
                if st.button('✨ AI 프레젠테이션 생성', type='primary', use_container_width=True):
                    with st.spinner('1/3  구조 에이전트 실행 중...'):
                        s_data = run_structure_agent(
                            st.session_state['topic'], lang,
                            st.session_state['category'],
                            st.session_state['slide_count'],
                            st.session_state['source_text'],
                        )
                    if not s_data:
                        st.error('구조 설계 실패. 다시 시도해 주세요.'); st.stop()
                    st.session_state['structure_data'] = s_data
                    st.session_state['agent_stage'] = 1
                    with st.spinner('2/3  콘텐츠 에이전트 실행 중...'):
                        c_data = run_content_agent(
                            s_data, st.session_state['topic'], lang,
                            st.session_state['category'], st.session_state['source_text'],
                        )
                    if not c_data:
                        st.error('콘텐츠 생성 실패. 다시 시도해 주세요.'); st.stop()
                    if use_polish:
                        with st.spinner('2.5/3  품질 에이전트 실행 중...'):
                            c_data = run_polish_agent(c_data, lang)
                    st.session_state['presentation_data'] = c_data
                    st.session_state['agent_stage'] = 2
                    with st.spinner('3/3  PPTX 컴파일 중...'):
                        trigger_pptx_recompile()
                    st.session_state['agent_stage'] = 3
                    st.rerun()

        else:
            # ── 생성 후: 문서 편집기 + AI 채팅 ───────────────────────
            doc_col, chat_col = st.columns([6, 4], gap='large')

            with doc_col:
                # 재생성 버튼
                rg_col, _ = st.columns([2, 6])
                with rg_col:
                    if st.button('🔄 다시 생성', use_container_width=True):
                        st.session_state['presentation_data'] = None
                        st.session_state['agent_stage'] = 0
                        st.session_state['chat_history'] = []
                        st.rerun()

                st.markdown('<div style="height:4px;"></div>', unsafe_allow_html=True)

                _TYPE_LABEL = {
                    'title_and_content': 'Content', 'section_header': 'Section',
                    'big_stat': 'Big Stat', 'two_column': 'Two Column',
                    'three_cards': 'Cards', 'timeline': 'Timeline', 'team_grid': 'Team',
                }

                def _esc(s):
                    return str(s or '').replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')

                def _doc_slide_html(sl, acc, txt):
                    """슬라이드 하나를 읽기 전용 문서 HTML로 렌더링."""
                    stype  = sl.get('slide_type', 'title_and_content')
                    title  = _esc(sl.get('title', ''))
                    buls   = [_esc(b) for b in sl.get('bullets', [])]
                    kw     = _esc(sl.get('key_takeaway', ''))

                    out = f'<div style="font-size:1.1rem;font-weight:700;color:{txt};margin-bottom:10px;line-height:1.3;">{title}</div>'

                    if stype == 'section_header':
                        sub = _esc(sl.get('summary', ''))
                        out += f'<div style="font-size:0.88rem;color:#888;margin-top:-4px;">{sub}</div>'

                    elif stype == 'big_stat':
                        sv = _esc(sl.get('stat_value', ''))
                        sd = _esc(sl.get('stat_description', ''))
                        out += (f'<div style="display:flex;align-items:baseline;gap:10px;margin-bottom:10px;">'
                                f'<span style="font-size:2.4rem;font-weight:900;color:{acc};line-height:1;">{sv}</span>'
                                f'<span style="font-size:0.85rem;color:#888;">{sd}</span></div>')
                        for b in buls:
                            out += f'<div style="font-size:0.88rem;color:#555;padding:2px 0 2px 12px;border-left:2px solid {acc}20;margin-bottom:4px;">• {b}</div>'

                    elif stype == 'three_cards':
                        cards = sl.get('cards', [])
                        out += '<div style="display:grid;grid-template-columns:repeat(3,1fr);gap:8px;margin-top:4px;">'
                        for c in cards[:3]:
                            ct = _esc(c.get('card_title','')); cc = _esc(c.get('card_content',''))
                            out += (f'<div style="background:{acc}08;border-radius:8px;padding:10px 12px;border:1px solid {acc}15;">'
                                    f'<div style="font-size:0.78rem;font-weight:700;color:{acc};margin-bottom:4px;">{ct}</div>'
                                    f'<div style="font-size:0.76rem;color:#666;line-height:1.5;">{cc}</div></div>')
                        out += '</div>'

                    elif stype == 'timeline':
                        steps_ = sl.get('timeline_steps', [])
                        for j, s_ in enumerate(steps_[:4]):
                            st_t = _esc(s_.get('step_title','')); st_d = _esc(s_.get('step_desc',''))
                            out += (f'<div style="display:flex;gap:10px;margin-bottom:8px;align-items:flex-start;">'
                                    f'<div style="width:20px;height:20px;border-radius:50%;background:{acc};flex-shrink:0;'
                                    f'display:flex;align-items:center;justify-content:center;'
                                    f'font-size:0.6rem;color:white;font-weight:800;">{j+1}</div>'
                                    f'<div><b style="font-size:0.86rem;color:{txt};">{st_t}</b>'
                                    f'<span style="font-size:0.8rem;color:#888;margin-left:8px;">{st_d}</span></div></div>')

                    elif stype == 'team_grid':
                        members = sl.get('team_members', [])
                        out += '<div style="display:flex;gap:14px;flex-wrap:wrap;margin-top:6px;">'
                        for mem in members:
                            n_ = _esc(mem.get('name','')); r_ = _esc(mem.get('role',''))
                            init = n_[0] if n_ else '?'
                            out += (f'<div style="text-align:center;min-width:60px;">'
                                    f'<div style="width:36px;height:36px;border-radius:50%;background:{acc};'
                                    f'margin:0 auto 5px;display:flex;align-items:center;justify-content:center;'
                                    f'font-size:0.9rem;color:white;font-weight:700;">{init}</div>'
                                    f'<div style="font-size:0.8rem;font-weight:600;color:{txt};">{n_}</div>'
                                    f'<div style="font-size:0.7rem;color:#aaa;">{r_}</div></div>')
                        out += '</div>'

                    else:  # title_and_content, two_column
                        half = max(len(buls)//2, 1) if stype == 'two_column' else len(buls)
                        if stype == 'two_column' and buls:
                            out += '<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;">'
                            for col_buls in [buls[:half], buls[half:]]:
                                out += '<div>'
                                for b in col_buls:
                                    out += f'<div style="font-size:0.88rem;color:#555;padding:3px 0;display:flex;gap:8px;"><span style="color:{acc};flex-shrink:0;">•</span>{b}</div>'
                                out += '</div>'
                            out += '</div>'
                        else:
                            for b in buls:
                                out += f'<div style="font-size:0.88rem;color:#555;padding:3px 0;display:flex;gap:8px;"><span style="color:{acc};flex-shrink:0;">•</span>{b}</div>'

                    if kw:
                        out += (f'<div style="margin-top:10px;padding:7px 12px;background:{acc}09;'
                                f'border-radius:6px;font-size:0.78rem;color:#777;border-left:3px solid {acc};">💡 {kw}</div>')
                    return out

                _LAYOUT_OPTIONS = [
                    ('title_and_content', 'Content'),
                    ('section_header',    'Section Header'),
                    ('big_stat',          'Big Stat'),
                    ('three_cards',       'Three Cards'),
                    ('timeline',          'Timeline'),
                    ('two_column',        'Two Column'),
                    ('team_grid',         'Team Grid'),
                ]
                _LAYOUT_KEYS   = [k for k,_ in _LAYOUT_OPTIONS]
                _LAYOUT_LABELS = [f"{l}" for _,l in _LAYOUT_OPTIONS]

                for i, slide in enumerate(slides):
                    sid   = slide.get('_id', str(i))
                    stype = slide.get('slide_type', 'title_and_content')

                    # 슬라이드 라벨 + 레이아웃 변경 + 삭제
                    hd_c, lt_c, del_c = st.columns([3, 5, 1])
                    with hd_c:
                        st.markdown(
                            f'<div style="font-size:0.6rem;font-weight:700;letter-spacing:0.14em;'
                            f'color:#c8c8dc;margin-top:22px;margin-bottom:6px;text-transform:uppercase;">'
                            f'SLIDE {i+1}</div>',
                            unsafe_allow_html=True
                        )
                    with lt_c:
                        cur_idx = _LAYOUT_KEYS.index(stype) if stype in _LAYOUT_KEYS else 0
                        new_type = st.selectbox(
                            '레이아웃',
                            options=_LAYOUT_KEYS,
                            format_func=lambda k: dict(_LAYOUT_OPTIONS).get(k, k),
                            index=cur_idx,
                            key=f'layout_{sid}',
                            label_visibility='collapsed',
                        )
                        if new_type != stype:
                            slide['slide_type'] = new_type
                            data['slides'][i] = slide
                            st.session_state['presentation_data'] = data
                            trigger_pptx_recompile()
                            st.rerun()
                    with del_c:
                        if st.button('✕', key=f'del_{sid}', help='삭제'):
                            slides.pop(i)
                            for idx2, s2 in enumerate(slides): s2['slide_index'] = idx2+1
                            data['slides'] = slides
                            st.session_state['presentation_data'] = data; st.rerun()

                    # ── PPT 캔버스 편집기 (contenteditable) ──────────────
                    updated = _slide_editor_component(
                        slide=slide,
                        acc=_s3_acc,
                        bg=_s3_bg,
                        txt=_s3_txt,
                        key=f'slide_editor_{sid}',
                    )
                    if updated and updated != slide:
                        updated['_id'] = sid
                        updated['slide_index'] = i + 1
                        data['slides'][i] = updated
                        st.session_state['presentation_data'] = data
                        trigger_pptx_recompile()

                    # 저장 & 재컴파일 버튼
                    if st.button('저장 & 재컴파일', key=f'save_{sid}', type='primary', use_container_width=True):
                        data['slides'][i] = st.session_state.get(f'slide_editor_{sid}') or slide
                        st.session_state['presentation_data'] = data
                        with st.spinner('재컴파일 중...'): trigger_pptx_recompile()
                        st.rerun()

                    st.markdown('<hr style="margin:10px 0;border-color:#f0f0f5;">', unsafe_allow_html=True)

                # 슬라이드 추가
                if st.button('➕ 슬라이드 추가', use_container_width=True):
                    blank = {
                        '_id': uuid.uuid4().hex[:8], 'slide_index': len(slides)+1,
                        'slide_type': 'title_and_content', 'title': '새 슬라이드',
                        'summary': '', 'bullets': ['내용을 입력하세요'], 'key_takeaway': '', 'speaker_notes': '',
                    }
                    data['slides'].append(blank)
                    st.session_state['presentation_data'] = data; st.rerun()

            # ── 오른쪽: AI 채팅 패널 ──────────────────────────────────
            with chat_col:
                # 채팅 히스토리 렌더링
                history_items = ''
                for msg in chat_history[-14:]:
                    role    = msg.get('role', 'user')
                    content = (msg.get('content', '')
                               .replace('&', '&amp;').replace('<', '&lt;')
                               .replace('>', '&gt;').replace('\n', '<br>'))
                    if role == 'user':
                        history_items += (
                            f'<div style="background:white;border:1px solid #e8e8f4;border-radius:10px;'
                            f'padding:10px 14px;margin-bottom:8px;font-size:0.84rem;color:#333;">'
                            f'<span style="font-size:0.65rem;font-weight:600;color:#bbb;'
                            f'letter-spacing:0.06em;display:block;margin-bottom:4px;">나</span>'
                            f'{content}</div>'
                        )
                    else:
                        history_items += (
                            f'<div style="background:{_s3_acc}0E;border-radius:10px;'
                            f'padding:10px 14px;margin-bottom:8px;font-size:0.84rem;color:#333;">'
                            f'<span style="font-size:0.65rem;font-weight:600;color:{_s3_acc};'
                            f'letter-spacing:0.06em;display:block;margin-bottom:4px;">AI</span>'
                            f'{content}</div>'
                        )

                if not chat_history:
                    _chat_body = ('<div style="font-size:0.82rem;color:#ccc;text-align:center;padding:24px 0;">'
                                  '오른쪽 채팅으로 AI에게<br>수정을 요청해보세요.</div>')
                else:
                    _chat_body = history_items

                st.markdown(
                    f'<div style="background:#f8f9ff;border-radius:14px;padding:18px 20px 14px;'
                    f'border:1px solid #eaeaf4;margin-bottom:14px;min-height:180px;">'
                    f'<div style="font-weight:700;font-size:0.9rem;color:#0F0F2A;margin-bottom:14px;">'
                    f'✨ AI 편집 어시스턴트</div>'
                    f'{_chat_body}'
                    f'</div>',
                    unsafe_allow_html=True
                )

                user_msg = st.text_area(
                    '수정 요청',
                    placeholder=(
                        '예) 슬라이드 3 제목을 더 강렬하게\n'
                        '예) 전체 톤을 더 포멀하게\n'
                        '예) 슬라이드 2 불릿을 5개로 늘려줘'
                    ),
                    height=110,
                    key='chat_input',
                    label_visibility='collapsed',
                )

                send_c, clr_c = st.columns([3, 1])
                with send_c:
                    send_clicked = st.button('수정 요청 →', type='primary', use_container_width=True, key='chat_send')
                with clr_c:
                    clear_clicked = st.button('초기화', use_container_width=True, key='chat_clear')

                if send_clicked and user_msg.strip():
                    chat_history.append({'role': 'user', 'content': user_msg.strip()})
                    with st.spinner('AI 수정 중...'):
                        updated_data, reply = run_chat_edit(data, user_msg.strip(), lang)
                    chat_history.append({'role': 'ai', 'content': reply})
                    st.session_state['chat_history'] = chat_history
                    if updated_data:
                        st.session_state['presentation_data'] = updated_data
                        with st.spinner('재컴파일 중...'):
                            trigger_pptx_recompile()
                    st.rerun()

                if clear_clicked:
                    st.session_state['chat_history'] = []
                    st.rerun()

                st.markdown(
                    f'<div style="margin-top:10px;font-size:0.74rem;color:#c8c8d8;line-height:1.8;">'
                    f'<b style="color:{_s3_acc};">슬라이드 번호</b> 포함 → 해당 슬라이드만 수정<br>'
                    f'"전체" / "모든" / "다" → 전체 슬라이드 수정'
                    f'</div>',
                    unsafe_allow_html=True
                )
