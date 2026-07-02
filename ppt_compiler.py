import re
import os as _os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as _qn
from pptx.enum.shapes import MSO_SHAPE

COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BG_LIGHT = RGBColor(248, 250, 252)

# Default Design Tokens (Nordic Tech)
DEFAULT_THEME = {
    "bg_dark": RGBColor(15, 23, 42),
    "bg_light": RGBColor(248, 250, 252),
    "primary_text": RGBColor(30, 41, 59),
    "secondary_text": RGBColor(71, 85, 105),
    "accent": RGBColor(79, 70, 229),
    "accent_light": RGBColor(129, 140, 248),
    "accent_bg": RGBColor(238, 242, 255),
    "font_title": "Malgun Gothic",
    "font_body": "Malgun Gothic"
}

THEMES = {
    "Nordic Tech (Ice White & Indigo)": DEFAULT_THEME,
    "Obsidian Dark (Slate Dark & Cyan)": {
        "bg_dark": RGBColor(15, 23, 42),
        "bg_light": RGBColor(30, 41, 59),
        "primary_text": RGBColor(255, 255, 255),
        "secondary_text": RGBColor(148, 163, 184),
        "accent": RGBColor(6, 182, 212),
        "accent_light": RGBColor(103, 232, 249),
        "accent_bg": RGBColor(21, 94, 117),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Emerald Business (Mint & Forest)": {
        "bg_dark": RGBColor(6, 78, 59),
        "bg_light": RGBColor(240, 253, 250),
        "primary_text": RGBColor(15, 23, 42),
        "secondary_text": RGBColor(55, 65, 81),
        "accent": RGBColor(5, 150, 105),
        "accent_light": RGBColor(110, 231, 183),
        "accent_bg": RGBColor(209, 250, 229),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Sunset Warm (Cream & Terracotta)": {
        "bg_dark": RGBColor(67, 20, 7),
        "bg_light": RGBColor(255, 251, 235),
        "primary_text": RGBColor(69, 26, 3),
        "secondary_text": RGBColor(120, 53, 4),
        "accent": RGBColor(217, 119, 6),
        "accent_light": RGBColor(252, 211, 77),
        "accent_bg": RGBColor(254, 243, 199),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Midnight Purple (Deep Violet & Lavender)": {
        "bg_dark": RGBColor(23, 13, 47),
        "bg_light": RGBColor(30, 20, 60),
        "primary_text": RGBColor(237, 233, 254),
        "secondary_text": RGBColor(167, 139, 250),
        "accent": RGBColor(139, 92, 246),
        "accent_light": RGBColor(196, 181, 253),
        "accent_bg": RGBColor(46, 16, 101),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Ocean Blue (Navy & Sky)": {
        "bg_dark": RGBColor(3, 18, 60),
        "bg_light": RGBColor(240, 249, 255),
        "primary_text": RGBColor(7, 23, 68),
        "secondary_text": RGBColor(3, 105, 161),
        "accent": RGBColor(2, 132, 199),
        "accent_light": RGBColor(56, 189, 248),
        "accent_bg": RGBColor(224, 242, 254),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Rose Gold (Ivory & Blush)": {
        "bg_dark": RGBColor(120, 53, 15),
        "bg_light": RGBColor(255, 251, 245),
        "primary_text": RGBColor(41, 28, 20),
        "secondary_text": RGBColor(120, 80, 55),
        "accent": RGBColor(194, 103, 39),
        "accent_light": RGBColor(251, 191, 36),
        "accent_bg": RGBColor(255, 237, 213),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Monochrome (Pure White & Black)": {
        "bg_dark": RGBColor(17, 17, 17),
        "bg_light": RGBColor(255, 255, 255),
        "primary_text": RGBColor(17, 17, 17),
        "secondary_text": RGBColor(80, 80, 80),
        "accent": RGBColor(17, 17, 17),
        "accent_light": RGBColor(100, 100, 100),
        "accent_bg": RGBColor(245, 245, 245),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Cherry Blossom (Soft Pink & Rose)": {
        "bg_dark": RGBColor(136, 19, 55),
        "bg_light": RGBColor(255, 241, 245),
        "primary_text": RGBColor(67, 20, 40),
        "secondary_text": RGBColor(159, 66, 95),
        "accent": RGBColor(225, 29, 72),
        "accent_light": RGBColor(251, 113, 133),
        "accent_bg": RGBColor(255, 228, 230),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
    "Custom": {
        "bg_dark": RGBColor(15, 23, 42),
        "bg_light": RGBColor(248, 250, 252),
        "primary_text": RGBColor(30, 41, 59),
        "secondary_text": RGBColor(71, 85, 105),
        "accent": RGBColor(79, 70, 229),
        "accent_light": RGBColor(129, 140, 248),
        "accent_bg": RGBColor(238, 242, 255),
        "font_title": "Malgun Gothic",
        "font_body": "Malgun Gothic"
    },
}

# Template layout styles — controls header, background, bullet, and takeaway rendering
TEMPLATE_STYLES = {
    "Modern": {
        "header_type": "text_underline",  # title text + short accent underline
        "title_size": 28,
        "slide_bg_key": "bg_light",
        "bullet_prefix": "•  ",
        "bullet_color_key": "secondary_text",
        "use_takeaway": True,
        "section_bg_key": "accent",
        "section_title_white": True,
        "title_slide_dark": True,
    },
    "Bold": {
        "header_type": "full_bar",        # full-width accent bar across top
        "title_size": 30,
        "slide_bg_key": "bg_light",
        "bullet_prefix": "▶  ",
        "bullet_color_key": "primary_text",
        "use_takeaway": True,
        "section_bg_key": "accent",
        "section_title_white": True,
        "title_slide_dark": True,
    },
    "Minimal": {
        "header_type": "accent_title",    # title in accent color, thin full-width line
        "title_size": 30,
        "slide_bg_key": "bg_light",
        "bullet_prefix": "—  ",
        "bullet_color_key": "primary_text",
        "use_takeaway": False,
        "section_bg_key": "accent_bg",
        "section_title_white": False,     # accent-colored text on light bg
        "title_slide_dark": False,        # light title slide
    },
    "Dark": {
        "header_type": "full_bar",
        "title_size": 28,
        "slide_bg_key": "bg_dark",
        "bullet_prefix": "◆  ",
        "bullet_color_key": "accent_light",
        "use_takeaway": True,
        "section_bg_key": "bg_dark",
        "section_title_white": False,     # accent-colored text on dark bg
        "title_slide_dark": True,
    },
}


def set_slide_background(slide, prs, color):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def add_header_styled(slide, title_text, colors, scfg, left, header_top, width, slide_w_emu, content_top):
    """
    Renders slide header according to the template style config (scfg).
    - text_underline: title + short accent underline (Modern)
    - full_bar: full-width accent bar from top to content_top, white title inside (Bold, Dark)
    - accent_title: title in accent color, thin full-width horizontal line below (Minimal)
    """
    header_type = scfg.get("header_type", "text_underline")
    title_size  = scfg.get("title_size", 28)

    if header_type == "full_bar":
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w_emu, content_top)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors["accent"]
        bar.line.fill.background()

        tb = slide.shapes.add_textbox(left, header_top, width, content_top - header_top)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_top = 0
        tf.margin_right = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = colors["font_title"]
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    elif header_type == "accent_title":
        tb = slide.shapes.add_textbox(left, header_top, width, Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = colors["font_title"]
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = colors["accent"]

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, header_top + Inches(0.85), slide_w_emu, Inches(0.025)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = colors["accent"]
        line.line.fill.background()

    else:  # text_underline (Modern)
        tb = slide.shapes.add_textbox(left, header_top, width, Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = colors["font_title"]
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = colors["primary_text"]

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, header_top + Inches(0.8), Inches(1.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = colors["accent"]
        line.line.fill.background()


def add_takeaway_banner(slide, text, colors, left, top, width, height):
    if not text:
        return

    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = colors["accent_bg"]
    banner.line.fill.background()

    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.3)
    tf.margin_bottom = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = f"💡 Key Takeaway: {text}"
    p.font.name = colors["font_body"]
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = colors["accent"]


def add_textbox_content(slide, text_list, colors, left, top, width, height,
                        font_size=15, bullet_prefix="•  ", bullet_color_key="secondary_text"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    first = True
    for bullet in text_list:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"{bullet_prefix}{bullet}"
        p.font.name = colors["font_body"]
        p.font.size = Pt(font_size)
        p.font.color.rgb = colors.get(bullet_color_key, colors["secondary_text"])
        p.space_after = Pt(10)
    return box


def compile_presentation(data, output_filename_or_stream, theme_name=None,
                         deck_format="Standard (16:9)", template_style="Modern"):
    colors = DEFAULT_THEME
    if theme_name and theme_name in THEMES:
        colors = THEMES[theme_name]

    scfg = TEMPLATE_STYLES.get(template_style, TEMPLATE_STYLES["Modern"])

    prs = Presentation()

    _FORMAT_DIMS = {
        "Standard (16:9)":                   (13.333, 7.5),
        "Keynote Widescreen (16:10)":         (13.333, 8.333),
        "Classic Slide (4:3)":               (10.0,   7.5),
        "Social Carousel (4:5)":             (8.0,    10.0),
        "Report / Document (A4-Horizontal)": (11.693, 8.268),
    }
    slide_w, slide_h = _FORMAT_DIMS.get(deck_format, (13.333, 7.5))
    prs.slide_width  = Inches(slide_w)
    prs.slide_height = Inches(slide_h)

    is_vertical = deck_format == "Social Carousel (4:5)"

    if is_vertical:
        margin_left    = Inches(0.8)
        content_width  = Inches(6.4)
        header_top     = Inches(0.6)
        content_top    = Inches(1.8)
        content_height = Inches(6.0)
        takeaway_top   = Inches(8.4)
        takeaway_height = Inches(1.0)
        title_top      = Inches(2.8)
        title_height   = Inches(4.0)
    else:
        margin_left    = Inches(slide_w * 0.075)
        content_width  = Inches(slide_w * 0.85)
        header_top     = Inches(slide_h * 0.08)
        content_top    = Inches(slide_h * 0.24)
        content_height = Inches(slide_h * 0.50)
        takeaway_top   = Inches(slide_h * 0.785)
        takeaway_height = Inches(slide_h * 0.107)
        title_top      = Inches(slide_h * 0.293)
        title_height   = Inches(slide_h * 0.533)

    blank_layout = prs.slide_layouts[6]

    # ── Title Slide ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)

    if scfg["title_slide_dark"]:
        set_slide_background(slide, prs, colors["bg_dark"])
        # Thin accent top bar decoration
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors["accent"]
        top_bar.line.fill.background()
        title_rgb = COLOR_WHITE
        sub_rgb   = colors["accent_light"]
    else:
        # Minimal style: light title slide
        set_slide_background(slide, prs, colors["bg_light"])
        # Large accent block on left (decorative)
        accent_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(slide_w * 0.05), prs.slide_height)
        accent_block.fill.solid()
        accent_block.fill.fore_color.rgb = colors["accent"]
        accent_block.line.fill.background()
        title_rgb = colors["primary_text"]
        sub_rgb   = colors["accent"]

    title_box = slide.shapes.add_textbox(margin_left, title_top, content_width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_title = tf.paragraphs[0]
    p_title.text = data.get("presentation_title", "AI Generated Presentation")
    p_title.font.name = colors["font_title"]
    p_title.font.size = Pt(38 if is_vertical else 44)
    p_title.font.bold = True
    p_title.font.color.rgb = title_rgb
    p_title.space_after = Pt(20)

    p_sub = tf.add_paragraph()
    p_sub.text = data.get("presentation_subtitle", "Generated by Snapdeck PPT AI Engine")
    p_sub.font.name = colors["font_body"]
    p_sub.font.size = Pt(16 if is_vertical else 18)
    p_sub.font.color.rgb = sub_rgb

    # ── Content Slides ───────────────────────────────────────────────────────
    slide_bg_color  = colors[scfg["slide_bg_key"]]
    bullet_prefix   = scfg["bullet_prefix"]
    bullet_clr_key  = scfg["bullet_color_key"]

    for item in data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        slide_type = item.get("slide_type", "title_and_content")

        if slide_type == "section_header":
            section_bg = colors[scfg["section_bg_key"]]
            set_slide_background(slide, prs, section_bg)

            if scfg["section_title_white"]:
                s_title_color = COLOR_WHITE
                s_sub_color   = colors["accent_bg"]
            else:
                # Accent text (works on both light accent_bg and dark bg_dark)
                s_title_color = colors["accent"]
                s_sub_color   = colors["secondary_text"]

            box = slide.shapes.add_textbox(
                margin_left,
                Inches(3.0) if is_vertical else Inches(2.5),
                content_width, Inches(3.5)
            )
            tf = box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = item.get("title", "")
            p.font.name = colors["font_title"]
            p.font.size = Pt(32 if is_vertical else 36)
            p.font.bold = True
            p.font.color.rgb = s_title_color
            p.space_after = Pt(14)

            p_sub = tf.add_paragraph()
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.text = item.get("summary", "")
            p_sub.font.name = colors["font_body"]
            p_sub.font.size = Pt(15 if is_vertical else 16)
            p_sub.font.color.rgb = s_sub_color

        elif slide_type == "big_stat":
            set_slide_background(slide, prs, slide_bg_color)
            add_header_styled(slide, item.get("title", ""), colors, scfg,
                              margin_left, header_top, content_width, prs.slide_width, content_top)

            if is_vertical:
                stat_box = slide.shapes.add_textbox(margin_left, Inches(2.2), content_width, Inches(2.5))
                tf_stat = stat_box.text_frame
                tf_stat.word_wrap = True

                p_val = tf_stat.paragraphs[0]
                p_val.text = item.get("stat_value", "100%")
                p_val.font.name = colors["font_title"]
                p_val.font.size = Pt(64)
                p_val.font.bold = True
                p_val.font.color.rgb = colors["accent"]

                p_desc = tf_stat.add_paragraph()
                p_desc.text = item.get("stat_description", "Metric Details")
                p_desc.font.name = colors["font_body"]
                p_desc.font.size = Pt(15)
                p_desc.font.bold = True
                p_desc.font.color.rgb = colors["primary_text"]

                content_box = slide.shapes.add_textbox(margin_left, Inches(4.5), content_width, Inches(3.6))
                tf_content = content_box.text_frame
            else:
                stat_box = slide.shapes.add_textbox(margin_left, Inches(2.2), Inches(4.5), Inches(3.2))
                tf_stat = stat_box.text_frame
                tf_stat.word_wrap = True

                p_val = tf_stat.paragraphs[0]
                p_val.text = item.get("stat_value", "100%")
                p_val.font.name = colors["font_title"]
                p_val.font.size = Pt(72)
                p_val.font.bold = True
                p_val.font.color.rgb = colors["accent"]
                p_val.space_after = Pt(10)

                p_desc = tf_stat.add_paragraph()
                p_desc.text = item.get("stat_description", "Metric Details")
                p_desc.font.name = colors["font_body"]
                p_desc.font.size = Pt(16)
                p_desc.font.bold = True
                p_desc.font.color.rgb = colors[bullet_clr_key]

                content_box = slide.shapes.add_textbox(Inches(6.0), Inches(2.2), Inches(6.33), Inches(3.2))
                tf_content = content_box.text_frame

            tf_content.word_wrap = True
            tf_content.margin_left = tf_content.margin_top = 0

            first = True
            for bullet in item.get("bullets", []):
                p = tf_content.paragraphs[0] if first else tf_content.add_paragraph()
                first = False
                p.text = f"{bullet_prefix}{bullet}"
                p.font.name = colors["font_body"]
                p.font.size = Pt(14 if is_vertical else 15)
                p.font.color.rgb = colors.get(bullet_clr_key, colors["secondary_text"])
                p.space_after = Pt(12)

            if scfg["use_takeaway"]:
                add_takeaway_banner(slide, item.get("key_takeaway", ""), colors,
                                    margin_left, takeaway_top, content_width, takeaway_height)

        elif slide_type == "two_column":
            set_slide_background(slide, prs, slide_bg_color)
            add_header_styled(slide, item.get("title", ""), colors, scfg,
                              margin_left, header_top, content_width, prs.slide_width, content_top)

            bullets = item.get("bullets", [])
            half = len(bullets) // 2
            left_bullets  = bullets[:half] if half > 0 else bullets
            right_bullets = bullets[half:] if half > 0 else []

            if is_vertical:
                add_textbox_content(slide, left_bullets, colors, margin_left, content_top,
                                    content_width, Inches(3.0), font_size=14,
                                    bullet_prefix=bullet_prefix, bullet_color_key=bullet_clr_key)
                if right_bullets:
                    add_textbox_content(slide, right_bullets, colors, margin_left,
                                        content_top + Inches(3.2), content_width, Inches(3.0),
                                        font_size=14, bullet_prefix=bullet_prefix,
                                        bullet_color_key=bullet_clr_key)
            else:
                col_w = (content_width - Inches(0.8)) / 2
                add_textbox_content(slide, left_bullets, colors, margin_left, content_top,
                                    col_w, content_height, font_size=15,
                                    bullet_prefix=bullet_prefix, bullet_color_key=bullet_clr_key)
                if right_bullets:
                    add_textbox_content(slide, right_bullets, colors,
                                        margin_left + col_w + Inches(0.8), content_top,
                                        col_w, content_height, font_size=15,
                                        bullet_prefix=bullet_prefix, bullet_color_key=bullet_clr_key)

            if scfg["use_takeaway"]:
                add_takeaway_banner(slide, item.get("key_takeaway", ""), colors,
                                    margin_left, takeaway_top, content_width, takeaway_height)

        elif slide_type == "three_cards":
            set_slide_background(slide, prs, slide_bg_color)
            add_header_styled(slide, item.get("title", ""), colors, scfg,
                              margin_left, header_top, content_width, prs.slide_width, content_top)

            cards = item.get("cards", [])
            if not cards:
                bullets = item.get("bullets", [])
                cards = [{"card_title": f"Point {idx+1}", "card_content": b}
                         for idx, b in enumerate(bullets[:3])]

            # Card background — for Dark style use accent_bg (which is accent tinted)
            card_fill = colors["accent_bg"]

            if is_vertical:
                card_h   = Inches(1.8)
                spacing  = Inches(0.2)
                for idx, card in enumerate(cards[:3]):
                    card_top = content_top + idx * (card_h + spacing)
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   margin_left, card_top, content_width, card_h)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = card_fill
                    shape.line.fill.background()

                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_top = Inches(0.15)

                    p_t = tf.paragraphs[0]
                    p_t.text = card.get("card_title", "")
                    p_t.font.name = colors["font_title"]
                    p_t.font.size = Pt(15)
                    p_t.font.bold = True
                    p_t.font.color.rgb = colors["accent"]
                    p_t.space_after = Pt(4)

                    p_c = tf.add_paragraph()
                    p_c.text = card.get("card_content", "")
                    p_c.font.name = colors["font_body"]
                    p_c.font.size = Pt(12)
                    p_c.font.color.rgb = colors[bullet_clr_key]
            else:
                card_w  = (content_width - Inches(0.8)) / 3
                spacing = Inches(0.4)
                for idx, card in enumerate(cards[:3]):
                    card_left = margin_left + idx * (card_w + spacing)
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   card_left, content_top, card_w,
                                                   content_height - Inches(0.2))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = card_fill
                    shape.line.fill.background()

                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_top = Inches(0.2)
                    tf.margin_right = Inches(0.2)

                    p_t = tf.paragraphs[0]
                    p_t.text = card.get("card_title", "")
                    p_t.font.name = colors["font_title"]
                    p_t.font.size = Pt(16)
                    p_t.font.bold = True
                    p_t.font.color.rgb = colors["accent"]
                    p_t.space_after = Pt(8)

                    p_c = tf.add_paragraph()
                    p_c.text = card.get("card_content", "")
                    p_c.font.name = colors["font_body"]
                    p_c.font.size = Pt(13)
                    p_c.font.color.rgb = colors[bullet_clr_key]

            if scfg["use_takeaway"]:
                add_takeaway_banner(slide, item.get("key_takeaway", ""), colors,
                                    margin_left, takeaway_top, content_width, takeaway_height)

        elif slide_type == "timeline":
            set_slide_background(slide, prs, slide_bg_color)
            add_header_styled(slide, item.get("title", ""), colors, scfg,
                              margin_left, header_top, content_width, prs.slide_width, content_top)

            steps = item.get("timeline_steps", [])
            if not steps:
                bullets = item.get("bullets", [])
                steps = [{"step_title": f"Phase {idx+1}", "step_desc": b}
                         for idx, b in enumerate(bullets[:4])]

            if is_vertical:
                line_left = margin_left + Inches(0.5)
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              line_left, content_top, Inches(0.04), Inches(5.8))
                line.fill.solid()
                line.fill.fore_color.rgb = colors["accent_light"]
                line.line.fill.background()

                step_h = Inches(1.3)
                for idx, step in enumerate(steps[:4]):
                    step_top = content_top + idx * Inches(1.4)

                    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                                    line_left - Inches(0.13), step_top + Inches(0.2),
                                                    Inches(0.3), Inches(0.3))
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = colors["accent"]
                    circle.line.fill.background()

                    box = slide.shapes.add_textbox(line_left + Inches(0.4), step_top,
                                                   content_width - Inches(0.9), step_h)
                    tf = box.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_top = 0

                    p_t = tf.paragraphs[0]
                    p_t.text = step.get("step_title", "")
                    p_t.font.name = colors["font_title"]
                    p_t.font.size = Pt(14)
                    p_t.font.bold = True
                    p_t.font.color.rgb = colors[bullet_clr_key]
                    p_t.space_after = Pt(2)

                    p_d = tf.add_paragraph()
                    p_d.text = step.get("step_desc", "")
                    p_d.font.name = colors["font_body"]
                    p_d.font.size = Pt(11)
                    p_d.font.color.rgb = colors.get(bullet_clr_key, colors["secondary_text"])
            else:
                line_top = content_top + Inches(1.3)
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              margin_left, line_top, content_width, Inches(0.04))
                line.fill.solid()
                line.fill.fore_color.rgb = colors["accent_light"]
                line.line.fill.background()

                num_steps = min(len(steps), 4)
                step_w = content_width / num_steps

                for idx, step in enumerate(steps[:num_steps]):
                    step_left = margin_left + idx * step_w

                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        step_left + (step_w / 2) - Inches(0.15), line_top - Inches(0.13),
                        Inches(0.3), Inches(0.3)
                    )
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = colors["accent"]
                    circle.line.fill.background()

                    is_above = idx % 2 == 0
                    box_top = line_top - Inches(1.1) if is_above else line_top + Inches(0.3)

                    box = slide.shapes.add_textbox(step_left + Inches(0.1), box_top,
                                                   step_w - Inches(0.2), Inches(1.0))
                    tf = box.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_top = 0

                    p_t = tf.paragraphs[0]
                    p_t.alignment = PP_ALIGN.CENTER
                    p_t.text = step.get("step_title", "")
                    p_t.font.name = colors["font_title"]
                    p_t.font.size = Pt(15)
                    p_t.font.bold = True
                    p_t.font.color.rgb = colors[bullet_clr_key]
                    p_t.space_after = Pt(4)

                    p_d = tf.add_paragraph()
                    p_d.alignment = PP_ALIGN.CENTER
                    p_d.text = step.get("step_desc", "")
                    p_d.font.name = colors["font_body"]
                    p_d.font.size = Pt(12)
                    p_d.font.color.rgb = colors.get(bullet_clr_key, colors["secondary_text"])

            if scfg["use_takeaway"]:
                add_takeaway_banner(slide, item.get("key_takeaway", ""), colors,
                                    margin_left, takeaway_top, content_width, takeaway_height)

        elif slide_type == "team_grid":
            set_slide_background(slide, prs, slide_bg_color)
            add_header_styled(slide, item.get("title", ""), colors, scfg,
                              margin_left, header_top, content_width, prs.slide_width, content_top)

            members = item.get("team_members", [])
            if not members:
                bullets = item.get("bullets", [])
                members = [
                    {
                        "name": b.split("-")[0].strip() if "-" in b else f"Member {idx+1}",
                        "role": "Lead" if idx == 0 else "Specialist",
                        "bio":  b.split("-")[1].strip() if "-" in b else b
                    }
                    for idx, b in enumerate(bullets[:3])
                ]

            # Card border color: accent_light for light themes, accent for dark
            card_border = colors["accent_light"]

            if is_vertical:
                card_h  = Inches(1.8)
                spacing = Inches(0.2)
                for idx, member in enumerate(members[:3]):
                    card_top = content_top + idx * (card_h + spacing)
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   margin_left, card_top, content_width, card_h)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = slide_bg_color
                    shape.line.color.rgb = card_border
                    shape.line.width = Pt(1)

                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_top = Inches(0.15)

                    p_n = tf.paragraphs[0]
                    p_n.text = member.get("name", "")
                    p_n.font.name = colors["font_title"]
                    p_n.font.size = Pt(15)
                    p_n.font.bold = True
                    p_n.font.color.rgb = colors[bullet_clr_key]
                    p_n.space_after = Pt(2)

                    p_r = tf.add_paragraph()
                    p_r.text = member.get("role", "")
                    p_r.font.name = colors["font_body"]
                    p_r.font.size = Pt(12)
                    p_r.font.bold = True
                    p_r.font.color.rgb = colors["accent"]
                    p_r.space_after = Pt(4)

                    p_b = tf.add_paragraph()
                    p_b.text = member.get("bio", "")
                    p_b.font.name = colors["font_body"]
                    p_b.font.size = Pt(11)
                    p_b.font.color.rgb = colors.get(bullet_clr_key, colors["secondary_text"])
            else:
                num_members = min(len(members), 3)
                card_w  = (content_width - Inches(0.8)) / num_members
                spacing = Inches(0.4)
                for idx, member in enumerate(members[:num_members]):
                    card_left = margin_left + idx * (card_w + spacing)
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   card_left, content_top, card_w,
                                                   content_height - Inches(0.2))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = slide_bg_color
                    shape.line.color.rgb = card_border
                    shape.line.width = Pt(1)

                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_top = Inches(0.2)
                    tf.margin_right = Inches(0.2)

                    p_n = tf.paragraphs[0]
                    p_n.text = member.get("name", "")
                    p_n.font.name = colors["font_title"]
                    p_n.font.size = Pt(16)
                    p_n.font.bold = True
                    p_n.font.color.rgb = colors[bullet_clr_key]
                    p_n.space_after = Pt(4)

                    p_r = tf.add_paragraph()
                    p_r.text = member.get("role", "")
                    p_r.font.name = colors["font_body"]
                    p_r.font.size = Pt(13)
                    p_r.font.bold = True
                    p_r.font.color.rgb = colors["accent"]
                    p_r.space_after = Pt(8)

                    p_b = tf.add_paragraph()
                    p_b.text = member.get("bio", "")
                    p_b.font.name = colors["font_body"]
                    p_b.font.size = Pt(11)
                    p_b.font.color.rgb = colors.get(bullet_clr_key, colors["secondary_text"])

            if scfg["use_takeaway"]:
                add_takeaway_banner(slide, item.get("key_takeaway", ""), colors,
                                    margin_left, takeaway_top, content_width, takeaway_height)

        else:  # title_and_content (default)
            set_slide_background(slide, prs, slide_bg_color)
            add_header_styled(slide, item.get("title", ""), colors, scfg,
                              margin_left, header_top, content_width, prs.slide_width, content_top)

            content_box = slide.shapes.add_textbox(margin_left, content_top, content_width, content_height)
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

            first = True
            for bullet in item.get("bullets", []):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = f"{bullet_prefix}{bullet}"
                p.font.name = colors["font_body"]
                p.font.size = Pt(15 if is_vertical else 16)
                p.font.color.rgb = colors.get(bullet_clr_key, colors["secondary_text"])
                p.space_after = Pt(14)

            if scfg["use_takeaway"]:
                add_takeaway_banner(slide, item.get("key_takeaway", ""), colors,
                                    margin_left, takeaway_top, content_width, takeaway_height)

    prs.save(output_filename_or_stream)


# ══════════════════════════════════════════════════════════════════
# TEMPLATE-BASED COMPILER  (방법 A)
# templates/ 폴더의 .pptx 파일을 슬라이드 디자인 베이스로 사용.
# 슬라이드 마스터·레이아웃의 배경·폰트·장식이 그대로 유지되며
# AI가 생성한 텍스트(제목, 불릿, 통계 등)가 삽입된다.
# ══════════════════════════════════════════════════════════════════

TEMPLATES_DIR = _os.path.join(_os.path.dirname(__file__), 'templates')


def list_templates():
    """templates/ 폴더에서 .pptx 파일 목록을 반환. [(표시이름, 경로), ...]"""
    if not _os.path.isdir(TEMPLATES_DIR):
        return []
    results = []
    for fname in sorted(_os.listdir(TEMPLATES_DIR)):
        if fname.lower().endswith('.pptx') and not fname.startswith('~$'):
            display = _os.path.splitext(fname)[0].replace('_', ' ').replace('-', ' ')
            results.append((display, _os.path.join(TEMPLATES_DIR, fname)))
    return results


def _clear_template_slides(prs):
    """템플릿의 콘텐츠 슬라이드를 모두 제거 (마스터/레이아웃은 보존)."""
    sld_id_lst = prs.slides._sldIdLst
    rIds = [el.get(_qn('r:id')) for el in list(sld_id_lst)]
    for el in list(sld_id_lst):
        sld_id_lst.remove(el)
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def _find_layout(prs, *keywords):
    """레이아웃 이름에 키워드가 포함된 첫 번째 레이아웃 반환."""
    for layout in prs.slide_layouts:
        name_lower = layout.name.lower()
        if any(k in name_lower for k in keywords):
            return layout
    return None


def _ph_fill(slide, idx, text, font_size=None, bold=False, color=None, align=None):
    """플레이스홀더 idx에 텍스트를 채움. 존재하면 True 반환."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx and ph.has_text_frame:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if align:
                p.alignment = align
            run = p.add_run()
            run.text = text
            if bold:
                run.font.bold = True
            if font_size:
                run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            return True
    return False


def _txb(slide, left_in, top_in, w_in, h_in, text,
         size=14, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    """슬라이드에 텍스트박스를 추가하고 반환."""
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def _add_bullets_txb(slide, left_in, top_in, w_in, h_in, bullets, size=15, color=None):
    """불릿 리스트를 텍스트박스로 추가."""
    if not bullets:
        return
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = '•  ' + b
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color
        p.space_after = Pt(10)


def _populate_template_slide(slide, slide_data, sw_in, sh_in):
    """
    템플릿 슬라이드에 AI 콘텐츠 삽입.
    - 플레이스홀더(idx=0 제목, idx=1 바디)가 있으면 우선 사용
    - 없으면 텍스트박스로 대체
    - 키테이크어웨이는 하단 텍스트박스로 항상 추가
    """
    stype   = slide_data.get('slide_type', 'title_and_content')
    title   = slide_data.get('title', '')
    summary = slide_data.get('summary', '')
    bullets = slide_data.get('bullets', [])
    kw      = slide_data.get('key_takeaway', '')

    # 기존 플레이스홀더 텍스트 초기화
    for ph in slide.placeholders:
        if ph.has_text_frame:
            try:
                ph.text_frame.clear()
            except Exception:
                pass

    # ── 제목 ──────────────────────────────────────────────────────────
    if not _ph_fill(slide, 0, title, font_size=28, bold=True):
        _txb(slide, 0.5, 0.25, sw_in - 1.0, 1.1, title, size=28, bold=True)

    # ── 콘텐츠 유형별 처리 ─────────────────────────────────────────────
    content_top  = 1.55
    content_h    = sh_in - content_top - 0.75
    content_w    = sw_in - 1.0
    kw_top       = sh_in - 0.65

    if stype == 'section_header':
        if not _ph_fill(slide, 1, summary, font_size=20):
            _txb(slide, 0.5, 1.7, content_w, 1.0, summary, size=20)
        return  # 섹션 헤더는 키테이크어웨이 없음

    elif stype in ('title_and_content', 'two_column'):
        if stype == 'two_column' and len(bullets) >= 2:
            half = len(bullets) // 2
            col_w = (content_w - 0.3) / 2
            _add_bullets_txb(slide, 0.5, content_top, col_w, content_h, bullets[:half])
            _add_bullets_txb(slide, 0.5 + col_w + 0.3, content_top, col_w, content_h, bullets[half:])
        else:
            if not _ph_fill(slide, 1, '\n'.join(bullets), font_size=15):
                _add_bullets_txb(slide, 0.5, content_top, content_w, content_h, bullets)

    elif stype == 'big_stat':
        stat_v = slide_data.get('stat_value', '')
        stat_d = slide_data.get('stat_description', '')
        _txb(slide, 0.5, content_top, content_w * 0.55, 1.6,
             stat_v, size=72, bold=True, align=PP_ALIGN.LEFT)
        _txb(slide, 0.5, content_top + 1.55, content_w, 0.5,
             stat_d, size=16)
        if bullets:
            _add_bullets_txb(slide, 0.5, content_top + 2.1, content_w, content_h - 2.1, bullets, size=14)

    elif stype == 'three_cards':
        cards = slide_data.get('cards', [])[:3]
        n = len(cards) or 1
        card_w = (content_w - 0.2 * (n - 1)) / n
        for k, card in enumerate(cards):
            left = 0.5 + k * (card_w + 0.2)
            _txb(slide, left, content_top, card_w, 0.45,
                 card.get('card_title', ''), size=14, bold=True)
            _txb(slide, left, content_top + 0.5, card_w, content_h - 0.55,
                 card.get('card_content', ''), size=12, wrap=True)

    elif stype == 'timeline':
        steps = slide_data.get('timeline_steps', [])[:4]
        row_h = min(content_h / max(len(steps), 1), 1.2)
        for k, step in enumerate(steps):
            top = content_top + k * row_h
            _txb(slide, 0.5, top, 0.4, row_h * 0.6,
                 str(k + 1), size=16, bold=True, align=PP_ALIGN.CENTER)
            _txb(slide, 1.0, top, content_w - 0.5, row_h * 0.45,
                 step.get('step_title', ''), size=14, bold=True)
            _txb(slide, 1.0, top + row_h * 0.44, content_w - 0.5, row_h * 0.5,
                 step.get('step_desc', ''), size=12)

    elif stype == 'team_grid':
        members = slide_data.get('team_members', [])[:8]
        cols = min(4, len(members) or 1)
        col_w = content_w / cols
        for k, mem in enumerate(members):
            col = k % cols
            row = k // cols
            left = 0.5 + col * col_w
            top  = content_top + row * 1.5
            _txb(slide, left, top, col_w - 0.1, 0.45,
                 mem.get('name', ''), size=13, bold=True, align=PP_ALIGN.CENTER)
            _txb(slide, left, top + 0.45, col_w - 0.1, 0.4,
                 mem.get('role', ''), size=11, align=PP_ALIGN.CENTER)

    # ── 키 테이크어웨이 (하단) ─────────────────────────────────────────
    if kw:
        _txb(slide, 0.5, kw_top, content_w, 0.55, f'💡  {kw}', size=11)


def compile_from_template(data, output_buf, template_path):
    """
    .pptx 템플릿 파일 기반 프레젠테이션 컴파일.

    템플릿의 슬라이드 마스터(배경, 색상, 폰트, 장식)를 보존하면서
    AI가 생성한 콘텐츠를 각 슬라이드에 삽입한다.
    """
    prs = Presentation(template_path)
    sw_in = prs.slide_width  / 914400
    sh_in = prs.slide_height / 914400

    # 1. 기존 콘텐츠 슬라이드 제거 (마스터/레이아웃 유지)
    _clear_template_slides(prs)

    # 2. 레이아웃 선택
    #    - 제목/섹션용  : "title", "cover", "intro", "divider", "section"
    #    - 콘텐츠용     : "content", "text", "body"
    #    - 빈 레이아웃  : "blank"
    layout_title = (
        _find_layout(prs, 'title slide', 'cover slide', 'intro', 'title_slide')
        or _find_layout(prs, 'title')
        or prs.slide_layouts[0]
    )
    layout_section = (
        _find_layout(prs, 'section', 'divider', 'chapter')
        or layout_title
    )
    layout_content = (
        _find_layout(prs, 'title and content', 'content', 'text and content')
        or _find_layout(prs, 'text', 'body')
        or (prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
    )
    layout_blank = (
        _find_layout(prs, 'blank')
        or prs.slide_layouts[-1]
    )

    # 3. 슬라이드 생성
    for i, slide_data in enumerate(data.get('slides', [])):
        stype = slide_data.get('slide_type', 'title_and_content')

        if i == 0:
            layout = layout_title
        elif stype == 'section_header':
            layout = layout_section
        elif stype in ('three_cards', 'big_stat', 'timeline', 'team_grid'):
            layout = layout_blank
        else:
            layout = layout_content

        slide = prs.slides.add_slide(layout)
        _populate_template_slide(slide, slide_data, sw_in, sh_in)

    prs.save(output_buf)
